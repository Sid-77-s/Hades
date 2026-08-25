from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_text, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title = slide.shapes.title
    title.text = title_text
    
    if layout_idx == 1:
        content = slide.shapes.placeholders[1]
        tf = content.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(content_text.split('\n')):
            if not line.strip():
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line.replace('- ', '')
            
            # Simple heuristic for bullets vs bold text
            if line.startswith('- '):
                p.level = 0
            elif line.startswith('  - '):
                p.level = 1
            elif line.startswith('    - '):
                p.level = 2
            else:
                p.level = 0
                
    return slide

prs = Presentation()

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "HADES\nThe AI Partner for Linux"
subtitle.text = "Open Innovation — Linux Based\n\nTurning natural-language goals into real work."

# Slide 2
add_slide(prs, "01 / The Problem", 
"""The AI Orchestration Burden
- Today's AI assistants are optimized for conversation, not execution.
- Real work is pushed back onto the user.
- The User is the Orchestrator:
  - You choose which model to use.
  - You decide which tool to open.
  - You copy/paste context between applications.
  - You execute commands manually.
  - You monitor progress and verify results.""")

# Slide 3
add_slide(prs, "02 / Why Current Workflows are Fragmented",
"""The Fragmentation Problem
- AI can explain what to do, but stops before the actual work.
- Complex tasks cross multiple tools:
  - Filesystem
  - Shell / Terminal
  - Browser
  - Research & Documents
- Users need an assistant that can act with context without turning every action into a permission dialog.""")

# Slide 4
add_slide(prs, "03 / The HADES Solution",
"""Outcome-Driven Execution
- HADES is an AI partner designed to turn natural-language goals into real work.
- Instead of micromanaging execution, the user gives an outcome.
- HADES handles the complexity:
  - Decides how to achieve the goal.
  - Delegates work to specialized models.
  - Uses Linux capabilities to execute.
  - Verifies the result.
  - Returns ONLY when the outcome is achieved.""")

# Slide 5
add_slide(prs, "04 / How HADES Works",
"""The Execution Loop
- 1. UNDERSTAND: Interpret the user's natural language intent.
- 2. PLAN: Break the outcome down into a structured mission.
- 3. DELEGATE: Assign specific steps to specialized workers.
- 4. EXECUTE: Use Linux tools (shell, files, browser) to perform the work.
- 5. VERIFY: Inspect results, identify failures, and retry if necessary.
- 6. RETURN: Report back to the user with a completed outcome.""")

# Slide 6
add_slide(prs, "05 / Architecture: The Brains",
"""Partner Brain + Execution Brain
- Partner Brain
  - Handles conversation and context.
  - Understands intent.
  - Decides if a request is casual chat or a real mission.
  - Clarifies ambiguity only when genuinely necessary.
- Execution Brain
  - Handles planning and orchestration.
  - Manages delegation to workers.
  - Monitors execution and handles results.""")

# Slide 7
add_slide(prs, "06 / Architecture: Workers & Tools",
"""Specialized Execution
- Workers
  - Replaceable, specialized AI models/agents.
  - Selected by HADES based on the task, not manually by the user.
- Linux Tools
  - Bash / Shell Execution
  - Filesystem Operations
  - Process Interaction
  - Browser & Web Research
  - APIs & Document Creation""")

# Slide 8
add_slide(prs, "07 / Novelty",
"""A Partner, Not a Command Box
- Partner Model
  - Takes ownership of the workflow instead of waiting for micro-instructions.
- Worker Architecture
  - Uses the right model for the right job, managed internally.
- Outcome-First Interaction
  - Instead of: "Run command A, open browser, save file..."
  - The User says: "Investigate this and give me the result."
- Linux Application Integration
  - Operates at the application layer with real system capabilities.""")

# Slide 9
add_slide(prs, "08 / Innovation",
"""Traditional AI vs. HADES
- Traditional AI:
  - USER -> PROMPT -> RESPONSE
- HADES:
  - USER -> GOAL -> HADES -> ORCHESTRATION -> WORKERS -> TOOLS -> VERIFIED RESULT
- The Technical Shift:
  - Human intent becomes an executable objective.
  - AI acts as an intelligent execution layer over Linux.
  - Focuses on verified completion, not just text generation.""")

# Slide 10
add_slide(prs, "09 / Real-World Use Case",
"""Intelligent Workspace Cleanup
- User: "Inspect my workspace, identify what is taking up space, and tell me what is safe to clean."
- HADES Execution:
  - Understands the goal.
  - Inspects the Linux environment.
  - Identifies large files (caches, logs, build artifacts).
  - Determines what is potentially safe to remove.
  - Does NOT blindly delete.
  - Returns a useful, verified report.
- Result: AI + Linux + Automation + Decision-making.""")

# Slide 11
add_slide(prs, "10 / Linux Integration",
"""Why Linux?
- The Target: An AI-powered application operating at the Linux application layer.
- Direct Control
  - Access to shell, filesystem, and processes provides a real execution environment.
- Composability
  - Integrates smoothly with open-source tooling, developer workflows, and scripting.
- Note: HADES is an AI-powered Linux application, NOT a kernel modification or OS replacement.""")

# Slide 12
add_slide(prs, "11 / The Demo",
"""One Complete End-to-End Mission
- The demo proves actual functionality, not simulated UI.
- The Flow:
  - 1. User gives HADES a real goal.
  - 2. HADES understands and creates a mission.
  - 3. Execution begins.
  - 4. Workers and Linux Tools are utilized.
  - 5. Result is verified.
  - 6. HADES returns to the user with the completed work.""")

# Slide 13
add_slide(prs, "12 / Dataset Used",
"""Foundation Models & System State
- Current Prototype
  - No custom training dataset is required.
  - Uses foundation models through configurable providers.
  - Relies on live system state, tool outputs, and application memory.
- Future Roadmap
  - Domain-specific evaluation datasets.
  - Task traces to improve reliability and execution speed.""")

# Slide 14
add_slide(prs, "13 / Impact & Future Vision",
"""AI That Works With You
- Current Impact
  - Reduces the cognitive load of orchestrating complex workflows.
  - Turns Linux into an outcome-driven environment.
- Future Vision (Roadmap)
  - Harden the Linux execution loop.
  - Strengthen verification and recovery strategies.
  - Expand worker capabilities and specialized agents.
  - Refine the partner-style UX and voice interaction.""")

# Slide 15
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "HADES"
subtitle.text = "Tell it what needs to happen.\nLet it handle the work.\n\nThank You."

prs.save('HADES_Pitch_Deck_Final.pptx')
print("Saved HADES_Pitch_Deck_Final.pptx")
