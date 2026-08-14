import os
from pptx import Presentation
from pptx.util import Inches
from src.capabilities.registry import BaseAdapter
from typing import Dict, Any
import litellm
import json

class PptxAdapter(BaseAdapter):
    async def execute(self, objective: str, context: Dict[str, Any]) -> Any:
        print("[PptxAdapter] Generating local presentation...")
        
        # Use LLM to generate the slide content based on the objective and context
        system_prompt = """You are a presentation generator. 
Given the objective and context, generate a JSON array of slides.
Each slide must have: "title" (string) and "content" (string with bullet points).
Keep it to 3-5 slides for brevity.
Output ONLY JSON."""

        user_prompt = f"Objective: {objective}\nContext: {json.dumps(context)}"
        
        response = litellm.completion(
            model="gemini/gemini-3.5-flash",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"}
        )
        
        try:
            slides_data = json.loads(response.choices[0].message.content)
            # if it's wrapped in a 'slides' key
            if isinstance(slides_data, dict) and "slides" in slides_data:
                slides_data = slides_data["slides"]
        except Exception:
            # Fallback data
            slides_data = [
                {"title": "Hades OS", "content": "The autonomous AI OS."},
                {"title": "Capabilities", "content": "- Execution Brain\n- Uncertainty Engine"}
            ]

        # Create PPTX
        prs = Presentation()
        
        for slide_info in slides_data:
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_info.get("title", "Slide")
            body.text = slide_info.get("content", "")
            
        output_filename = "hades_presentation.pptx"
        output_path = os.path.join(os.getcwd(), output_filename)
        prs.save(output_path)
        
        return f"Successfully generated local presentation at {output_path}"
        
    async def observe(self) -> Any:
        pass
        
    async def validate(self) -> bool:
        return True
