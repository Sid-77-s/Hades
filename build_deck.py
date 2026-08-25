"""
HackOHades.pptx — 12-slide HADES presentation
Careful pixel-perfect layout, no overlapping, dark theme.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PALETTE ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0a, 0x14)
PANEL   = RGBColor(0x12, 0x12, 0x28)
CARD    = RGBColor(0x18, 0x18, 0x35)
ION     = RGBColor(0x7c, 0x3a, 0xed)
ION2    = RGBColor(0x5b, 0x21, 0xb6)
GLOW    = RGBColor(0xa7, 0x8b, 0xfa)
WHITE   = RGBColor(0xff, 0xff, 0xff)
GREEN   = RGBColor(0x10, 0xb9, 0x81)
AMBER   = RGBColor(0xf5, 0x9e, 0x0b)
RED     = RGBColor(0xef, 0x44, 0x44)
SLATE   = RGBColor(0x94, 0xa3, 0xb8)
MUTED   = RGBColor(0x64, 0x74, 0x8b)
DKGREEN = RGBColor(0x04, 0x1e, 0x10)
DKAMBER = RGBColor(0x1c, 0x14, 0x02)
DKRED   = RGBColor(0x1c, 0x05, 0x05)
DKVIOLET= RGBColor(0x14, 0x08, 0x38)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── PRIMITIVES ────────────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(s, x, y, w, h, fill=PANEL, border=None, bw=Pt(1)):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = bw
    else:
        sh.line.fill.background()
    return sh

def text(s, txt, x, y, w, h, size=12, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def hline(s, x, y, w, color=ION, thickness=Pt(1.5)):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def badge(s, lbl, x, y, w, h, fg, bg):
    rect(s, x, y, w, h, fill=bg, border=fg, bw=Pt(0.75))
    text(s, lbl, x+0.05, y+0.04, w-0.1, h-0.08, size=8, bold=True,
         color=fg, align=PP_ALIGN.CENTER)

def card_box(s, x, y, w, h, title, lines, tc=GLOW, bc=ION,
             ts=11, bs=9, title_h=0.32):
    rect(s, x, y, w, h, fill=CARD, border=bc, bw=Pt(0.8))
    text(s, title, x+0.12, y+0.06, w-0.24, title_h,
         size=ts, bold=True, color=tc)
    body = "\n".join(lines)
    text(s, body, x+0.12, y+0.06+title_h+0.04, w-0.24,
         h-0.06-title_h-0.1, size=bs, color=SLATE, wrap=True)

def slide_num(s, n, total=12):
    text(s, f"{n:02d} / {total}", 12.5, 7.15, 0.7, 0.28,
         size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def top_label(s, txt_val):
    badge(s, txt_val, 0.3, 0.12, 12.73, 0.28, GLOW, DKVIOLET)

def section_title(s, main, sub=None):
    text(s, main, 0.3, 0.5, 12.73, 0.65,
         size=30, bold=True, color=WHITE)
    if sub:
        text(s, sub, 0.3, 1.12, 12.73, 0.38,
             size=15, color=GLOW)
    y = 1.55 if sub else 1.2
    hline(s, 0.3, y, 12.73)
    return y + 0.15


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = slide()

# Top strip
rect(s1, 0, 0, 13.33, 0.08, fill=ION, border=None)

# Central badge
badge(s1, "HACKATHON BUILD  ·  AI SYSTEMS ENGINEERING  ·  2025",
      3.0, 0.28, 7.33, 0.3, GLOW, DKVIOLET)

# HADES
text(s1, "H A D E S", 0, 1.1, 13.33, 2.2,
     size=108, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Violet underline
hline(s1, 1.5, 3.2, 10.33, ION, Pt(2.5))

# Subtitle
text(s1, "FROM AI ASSISTANT  →  AI OPERATING SYSTEM",
     0, 3.35, 13.33, 0.52,
     size=16, color=SLATE, align=PP_ALIGN.CENTER)

# Pipeline row
nodes = [
    ("HUMAN",                  MUTED,  PANEL),
    ("  →  ",                  ION,    None),
    ("HADES",                  WHITE,  ION2),
    ("  →  ",                  ION,    None),
    ("MODELS + TOOLS + SYSTEM",GLOW,   CARD),
    ("  →  ",                  ION,    None),
    ("REAL RESULT",            GREEN,  DKGREEN),
]
widths = [1.5, 0.55, 1.5, 0.55, 3.0, 0.55, 1.7]
x = 0.49
y_row = 4.05
for (label, fc, bg), w in zip(nodes, widths):
    if bg:
        rect(s1, x, y_row, w, 0.5, fill=bg, border=fc, bw=Pt(0.8))
    text(s1, label, x+0.05, y_row+0.1, w-0.1, 0.32,
         size=11, bold=True, color=fc, align=PP_ALIGN.CENTER)
    x += w + 0.03

# Bottom statement
hline(s1, 1.5, 4.82, 10.33, MUTED, Pt(0.5))
text(s1,
     "The user should manage the goal — not the intelligence required to achieve it.",
     0, 4.97, 13.33, 0.42,
     size=13, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# Tagline row
text(s1, "Linux-Native  ·  Model-Agnostic  ·  Mission-Locked  ·  Verified Execution",
     0, 5.55, 13.33, 0.35,
     size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom strip
rect(s1, 0, 7.38, 13.33, 0.12, fill=ION2, border=None)

slide_num(s1, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s2 = slide()
top_label(s2, "PROBLEM STATEMENT")
y0 = section_title(s2, "The AI Problem Is Not Capability.",
                   "It Is Orchestration.")

# Two columns
col_l = 0.3
col_r = 6.85
col_w = 6.3

# Left — fragmentation
rect(s2, col_l, y0, col_w, 0.32, fill=DKVIOLET, border=ION)
text(s2, "TODAY'S REALITY — YOU ARE THE ORCHESTRATOR",
     col_l+0.1, y0+0.06, col_w-0.2, 0.22,
     size=9, bold=True, color=GLOW)

tools = [
    ("🔍 Gemini",  "Research"),
    ("💻 Claude",  "Coding"),
    ("✍  ChatGPT", "Writing"),
    ("📝 Notion",  "Notes"),
    ("⚡ Zapier",  "Automation"),
    ("📅 Calendar","Scheduling"),
]
tw = 2.95; th = 0.48
positions = [
    (col_l,       y0+0.42),
    (col_l+3.05,  y0+0.42),
    (col_l,       y0+1.00),
    (col_l+3.05,  y0+1.00),
    (col_l,       y0+1.58),
    (col_l+3.05,  y0+1.58),
]
for (name, task), (tx, ty) in zip(tools, positions):
    rect(s2, tx, ty, tw, th, fill=CARD, border=MUTED, bw=Pt(0.6))
    text(s2, name, tx+0.1, ty+0.04, tw-0.2, 0.25, size=11, bold=True, color=WHITE)
    text(s2, task, tx+0.1, ty+0.27, tw-0.2, 0.18, size=9, color=SLATE)

# YOU box
rect(s2, col_l+2.15, y0+1.0, 1.7, 0.46, fill=RGBColor(0x25,0x05,0x05), border=RED)
text(s2, "YOU", col_l+2.15, y0+1.04, 1.7, 0.38,
     size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Context lost note
rect(s2, col_l, y0+2.18, col_w, 0.32, fill=DKRED, border=RED, bw=Pt(0.6))
text(s2, "Context is lost between every session. Every model starts fresh. Orchestration belongs to you.",
     col_l+0.1, y0+2.24, col_w-0.2, 0.22, size=9, color=RED)

# Right — hidden tax
rect(s2, col_r, y0, col_w, 0.32, fill=DKVIOLET, border=ION)
text(s2, "THE HIDDEN COGNITIVE TAX",
     col_r+0.1, y0+0.06, col_w-0.2, 0.22, size=9, bold=True, color=GLOW)

problems = [
    ("Context switching between tools",       "Every task requires switching mental context"),
    ("Repeating yourself to every model",      "No shared memory across AI systems"),
    ("Manually selecting the right AI",        "Which model? For what? At what cost?"),
    ("Managing workflow sequencing",           "You plan, you delegate, you track"),
    ("Context lost between sessions",          "Yesterday's work is invisible today"),
    ("Decision fatigue from tool sprawl",      "10+ tools = 10x the cognitive overhead"),
]
for i, (title, detail) in enumerate(problems):
    py = y0 + 0.42 + i * 0.6
    rect(s2, col_r, py, col_w, 0.5, fill=DKRED, border=RED, bw=Pt(0.6))
    text(s2, title,  col_r+0.12, py+0.04, col_w-0.24, 0.24, size=11, bold=True, color=WHITE)
    text(s2, detail, col_r+0.12, py+0.27, col_w-0.24, 0.18, size=9,  color=SLATE)

# Bottom
hline(s2, 0.3, 6.9, 12.73, MUTED, Pt(0.5))
text(s2,
     "Today's AI gives humans MORE capabilities.  HADES removes the burden of coordinating them.",
     0, 7.02, 13.33, 0.35, size=12, bold=True, color=GLOW, align=PP_ALIGN.CENTER)
slide_num(s2, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIVE
# ══════════════════════════════════════════════════════════════════════════════
s3 = slide()
top_label(s3, "OBJECTIVE")
y0 = section_title(s3, "What HADES Sets Out to Achieve.",
                   "A single AI layer that orchestrates everything, so you don't have to.")

# Large mission statement
rect(s3, 0.3, y0, 12.73, 0.85, fill=DKVIOLET, border=ION, bw=Pt(1.5))
text(s3,
     "Build an AI Operating System that converts human intent into verified, "
     "system-level action — through conversation, structured mission planning, "
     "multi-model orchestration, and outcome validation.",
     0.5, y0+0.1, 12.33, 0.65, size=14, color=WHITE, wrap=True)

y0 += 1.05

# 3x2 objective grid
objs = [
    ("01  ELIMINATE ORCHESTRATION BURDEN",
     "The user expresses a goal. HADES determines which models, tools, "
     "and sequences are needed — invisibly."),
    ("02  CONVERSATION BEFORE EXECUTION",
     "No AI action without mutual understanding. HADES aligns on objective, "
     "outcome, and constraints before touching the system."),
    ("03  MISSION-CENTRIC COMPUTING",
     "Replace disposable prompts with persistent Mission objects — structured, "
     "stateful, and validated end-to-end."),
    ("04  MODEL AGNOSTICISM",
     "No vendor lock-in. Gemini, Groq, OpenRouter, Ollama — all interchangeable "
     "via a unified LiteLLM abstraction layer."),
    ("05  VERIFIED EXECUTION",
     "Workers generate output. The Review Engine validates it. "
     "Generation is not success — the mission criteria define success."),
    ("06  LINUX-NATIVE OPERATION",
     "Terminal execution, filesystem control, process management — "
     "all running natively on the host Linux system with local TTS and local AI."),
]
cw = 4.17; ch = 1.12; gap = 0.1
for i, (title, body) in enumerate(objs):
    col = i % 3; row = i // 3
    ox = 0.3 + col * (cw + gap)
    oy = y0 + row * (ch + gap)
    rect(s3, ox, oy, cw, ch, fill=CARD, border=ION, bw=Pt(0.7))
    text(s3, title, ox+0.12, oy+0.07, cw-0.24, 0.3,
         size=9, bold=True, color=GLOW)
    text(s3, body,  ox+0.12, oy+0.40, cw-0.24, ch-0.52,
         size=9, color=SLATE, wrap=True)

slide_num(s3, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MARKET GAPS
# ══════════════════════════════════════════════════════════════════════════════
s4 = slide()
top_label(s4, "MARKET GAPS")
y0 = section_title(s4, "Where Existing AI Systems Fall Short.",
                   "11-dimension comparison across Traditional AI, AI Agents, and HADES.")

# Table
rows = [
    ("Conversation",              "✓", "~", "✓"),
    ("Persistent Identity",       "✗", "✗", "✓"),
    ("Long-Term Memory",          "✗", "~", "~"),
    ("Goal Understanding",        "~", "~", "✓"),
    ("Human Alignment Gate",      "✗", "✗", "✓"),
    ("Multi-Model Orchestration", "✗", "~", "✓"),
    ("Background Execution",      "✗", "✓", "✓"),
    ("Output Validation",         "✗", "✗", "✓"),
    ("Human Goal Authority",      "✗", "✗", "✓"),
    ("Continuous Relationship",   "✗", "✗", "~"),
    ("Local / Offline AI",        "✗", "~", "✓"),
]
col_x  = [0.3, 5.4, 8.2, 10.85]
col_w2 = [5.0, 2.7, 2.55, 2.35]
headers = ["CAPABILITY", "TRADITIONAL AI", "AI AGENT", "HADES"]
hcolors = [GLOW, SLATE, SLATE, WHITE]
hbg     = [DKVIOLET, CARD, CARD, ION2]
row_h   = 0.43

# Header
for ci, (hdr, hc, hbg_c, cw2) in enumerate(zip(headers, hcolors, hbg, col_w2)):
    rect(s4, col_x[ci], y0, cw2, 0.38, fill=hbg_c, border=ION, bw=Pt(0.8))
    text(s4, hdr, col_x[ci]+0.08, y0+0.07, cw2-0.16, 0.25,
         size=10, bold=True, color=hc, align=PP_ALIGN.CENTER)

y0 += 0.43
for ri, (label_t, v1, v2, v3) in enumerate(rows):
    bg_row = PANEL if ri % 2 == 0 else BG
    vals = [label_t, v1, v2, v3]
    for ci, (val, cw2) in enumerate(zip(vals, col_w2)):
        if ci == 0:
            rect(s4, col_x[ci], y0, cw2, row_h-0.02, fill=bg_row,
                 border=RGBColor(0x25,0x25,0x45), bw=Pt(0.5))
            text(s4, val, col_x[ci]+0.1, y0+0.1, cw2-0.2, row_h-0.14,
                 size=10, color=WHITE)
        else:
            v_map = {"✓": (GREEN, DKGREEN), "~": (AMBER, DKAMBER), "✗": (RED, DKRED)}
            fc, bc2 = v_map[val]
            hades_border = ION if ci == 3 else RGBColor(0x25,0x25,0x45)
            rect(s4, col_x[ci], y0, cw2, row_h-0.02,
                 fill=bc2, border=hades_border, bw=Pt(0.5 if ci < 3 else 1.0))
            text(s4, val, col_x[ci], y0+0.06, cw2, row_h-0.14,
                 size=14, bold=True, color=fc, align=PP_ALIGN.CENTER)
    y0 += row_h

# Bottom note
rect(s4, 0.3, y0+0.08, 12.73, 0.38, fill=DKVIOLET, border=ION, bw=Pt(0.8))
text(s4,
     "✓ = Implemented   ~= Partial   ✗ = Missing   |   HADES targets the full-row solution no existing system achieves.",
     0.45, y0+0.15, 12.5, 0.25, size=9, color=GLOW)

slide_num(s4, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM DESIGN
# ══════════════════════════════════════════════════════════════════════════════
s5 = slide()
top_label(s5, "SYSTEM DESIGN")
y0 = section_title(s5, "One Partner. Multiple Brains. Many Capabilities.",
                   "Layered AI OS architecture with strict separation of conversation and execution.")

layers = [
    # (label, sublabel, badge, badge_color, bg, border, note_color, note)
    ("USER",
     "Goal  ·  Context  ·  Approval",
     "HUMAN IN CONTROL", SLATE, PANEL, MUTED,
     None, None),

    ("PARTNER BRAIN",
     "Intent Classification  |  Mission Extraction  |  Understanding Evaluation  |  Conversational Alignment",
     "IMPLEMENTED ✓", GREEN, DKGREEN, GREEN,
     None, "partner_brain.py · intent_classifier.py · understanding_evaluator.py"),

    ("🔒  MISSION LOCK",
     "MUTUAL UNDERSTANDING GATE — Execution is gated until all mission fields + mutual_understanding_reached = True",
     "IMPLEMENTED ✓", WHITE, ION2, ION,
     None, "understanding_evaluator.py · mission.py → MissionStatus.AUTHORIZED_EXECUTION"),

    ("EXECUTIVE BRAIN",
     "Task Graph Generation  |  Worker Delegation  |  Retry & Recovery",
     "PARTIAL ⚡", AMBER, DKAMBER, AMBER,
     AMBER, "execution_brain.py  —  currently 1-task graph; dynamic decomposition = next milestone"),

    ("WORKER MANAGER  +  MODEL PROVIDERS",
     "Gemini (Google)  |  Llama 3.3 / Groq  |  OpenRouter  |  Ollama (local/offline)  — via LiteLLM with auto-fallback",
     "IMPLEMENTED ✓", GREEN, DKGREEN, GREEN,
     None, "worker_manager.py · config.json · LiteLLM"),

    ("TOOLS / SKILLS",
     "Terminal (bash)  |  Filesystem  |  Process Manager  |  Browser (partial)",
     "PARTIAL ⚡", AMBER, DKAMBER, AMBER,
     None, "skills/computer/terminal.py · filesystem.py · process_manager.py"),

    ("REVIEW ENGINE",
     "Exit Code Check  |  Artifact Existence  |  Heuristic Validation  |  LLM Semantic Review",
     "IMPLEMENTED ✓", GREEN, DKGREEN, GREEN,
     None, "review_engine.py — 4-stage validation before delivery"),

    ("MEMORY  +  SSE  +  TTS",
     "Session State (in-memory)  |  Mission History (JSON)  |  SSE event stream  |  Kokoro-ONNX local TTS",
     "PARTIAL ⚡", AMBER, DKAMBER, AMBER,
     None, "memory_manager.py · main.py /api/events · voice_manager.py"),
]

lh = 0.56
y = y0
for (lbl, sub, badge_t, badge_c, bg_c, border_c, note_c, file_ref) in layers:
    rect(s5, 0.3, y, 12.73, lh, fill=bg_c, border=border_c, bw=Pt(0.9))
    text(s5, lbl,  0.45, y+0.05,  5.5,  0.25, size=12, bold=True, color=badge_c)
    text(s5, sub,  0.45, y+0.30,  9.2,  0.20, size=8.5, color=SLATE)
    if file_ref:
        text(s5, file_ref, 0.45, y+0.42, 9.2, 0.12, size=7, color=MUTED, italic=True)
    # Badge
    bw2 = 2.0
    rect(s5, 11.0, y+0.13, bw2, 0.3, fill=bg_c, border=badge_c, bw=Pt(0.7))
    text(s5, badge_t, 11.05, y+0.17, bw2-0.1, 0.22,
         size=8, bold=True, color=badge_c, align=PP_ALIGN.CENTER)
    # Arrow
    if lbl != layers[-1][0]:
        text(s5, "▼", 6.5, y+lh, 0.35, 0.2, size=7, color=ION, align=PP_ALIGN.CENTER)
    y += lh + 0.2

slide_num(s5, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s6 = slide()
top_label(s6, "ARCHITECTURE DIAGRAM")
y0 = section_title(s6, "End-to-End Data Flow.",
                   "From user input → mission lock → execution → review → delivery.")

# Main flow in 3-column layout:
# LEFT: Input layer | CENTER: Processing | RIGHT: Output/Infra

# Column headers
for cx, cw2, ctitle, cc in [
    (0.3,  3.9, "INPUT LAYER",       GLOW),
    (4.45, 5.0, "PROCESSING CORE",   WHITE),
    (9.7,  3.4, "OUTPUT & INFRA",    GREEN),
]:
    rect(s6, cx, y0, cw2, 0.32, fill=DKVIOLET, border=ION, bw=Pt(0.8))
    text(s6, ctitle, cx+0.1, y0+0.06, cw2-0.2, 0.22,
         size=10, bold=True, color=cc, align=PP_ALIGN.CENTER)

y0 += 0.42

# LEFT COLUMN — input
left_items = [
    ("BROWSER UI",      "React + Vite + Tailwind",   "frontend/src/", ION),
    ("VOICE INPUT",     "Web Speech API (STT)",      "Browser native", AMBER),
    ("POST /api/chat",  "HTTP REST endpoint",         "main.py",        ION),
    ("SESSION STATE",   "UUID → SessionState dict",   "main.py",        SLATE),
]
for i, (t, d, f, c) in enumerate(left_items):
    yl = y0 + i * 1.02
    rect(s6, 0.3, yl, 3.9, 0.88, fill=CARD, border=c, bw=Pt(0.7))
    text(s6, t, 0.42, yl+0.04, 3.66, 0.28, size=11, bold=True, color=c)
    text(s6, d, 0.42, yl+0.30, 3.66, 0.22, size=9,  color=SLATE)
    text(s6, f, 0.42, yl+0.52, 3.66, 0.18, size=7,  color=MUTED, italic=True)
    if i < len(left_items)-1:
        text(s6, "▼", 2.1, yl+0.88, 0.5, 0.14, size=8, color=ION, align=PP_ALIGN.CENTER)

# CENTER COLUMN — core
center_items = [
    ("PARTNER BRAIN",       "Intent classify → extract mission → evaluate understanding",
     "partner_brain.py", GREEN),
    ("🔒 MISSION LOCK",     "AUTHORIZED_EXECUTION state · mutual_understanding_reached=True",
     "understanding_evaluator.py", WHITE),
    ("EXECUTIVE BRAIN",     "TaskGraph (1 task) → WorkerManager → LiteLLM routing",
     "execution_brain.py", AMBER),
    ("TOOLS",               "TerminalSkill  |  FilesystemSkill  |  ProcessManager",
     "skills/computer/", GREEN),
]
for i, (t, d, f, c) in enumerate(center_items):
    yl = y0 + i * 1.02
    bg_c2 = ION2 if t == "🔒 MISSION LOCK" else CARD
    rect(s6, 4.45, yl, 5.0, 0.88, fill=bg_c2, border=c, bw=Pt(0.9 if t=="🔒 MISSION LOCK" else 0.7))
    text(s6, t, 4.6, yl+0.04, 4.7, 0.28, size=11, bold=True, color=c)
    text(s6, d, 4.6, yl+0.30, 4.7, 0.22, size=9,  color=SLATE, wrap=True)
    text(s6, f, 4.6, yl+0.52, 4.7, 0.18, size=7,  color=MUTED, italic=True)
    if i < len(center_items)-1:
        text(s6, "▼", 6.85, yl+0.88, 0.5, 0.14, size=8, color=ION, align=PP_ALIGN.CENTER)

# RIGHT COLUMN — output
right_items = [
    ("REVIEW ENGINE",     "Exit code · Artifact · Heuristic · LLM semantic",
     "review_engine.py", GREEN),
    ("MEMORY",            "Session (RAM) · Mission history (JSON)",
     "memory_manager.py", AMBER),
    ("SSE EVENTS",        "Real-time stream to UI\n/api/events endpoint",
     "FastAPI EventSource", ION),
    ("TTS OUTPUT",        "Kokoro-ONNX → Base64 audio",
     "voice_manager.py", GREEN),
]
for i, (t, d, f, c) in enumerate(right_items):
    yl = y0 + i * 1.02
    rect(s6, 9.7, yl, 3.4, 0.88, fill=CARD, border=c, bw=Pt(0.7))
    text(s6, t, 9.82, yl+0.04, 3.16, 0.28, size=11, bold=True, color=c)
    text(s6, d, 9.82, yl+0.30, 3.16, 0.22, size=9,  color=SLATE, wrap=True)
    text(s6, f, 9.82, yl+0.52, 3.16, 0.18, size=7,  color=MUTED, italic=True)
    if i < len(right_items)-1:
        text(s6, "▼", 11.3, yl+0.88, 0.5, 0.14, size=8, color=ION, align=PP_ALIGN.CENTER)

# Horizontal connector arrows
for yl_a in [y0+0.3, y0+1.32, y0+2.34, y0+3.36]:
    text(s6, "→", 4.22, yl_a, 0.3, 0.28, size=13, color=GLOW, align=PP_ALIGN.CENTER)
    text(s6, "→", 9.45, yl_a, 0.3, 0.28, size=13, color=GLOW, align=PP_ALIGN.CENTER)

slide_num(s6, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MISSION LOCK
# ══════════════════════════════════════════════════════════════════════════════
s7 = slide()
top_label(s7, "HADES  —  MISSION LOCK")
y0 = section_title(s7, "HADES Doesn't Start With a Task.",
                   "It starts with understanding. Execution is gated by mutual alignment.")

# Left: pipeline
steps = [
    ("① USER INTENT",         SLATE,  PANEL),
    ("② CONVERSATION",         GREEN,  DKGREEN),
    ("③ CLARIFICATION",        GREEN,  DKGREEN),
    ("④ PUSHBACK / ALIGN",     GREEN,  DKGREEN),
    ("⑤ MUTUAL UNDERSTANDING", GREEN,  DKGREEN),
    ("🔒  MISSION LOCK",       WHITE,  ION2),
    ("⑦ PLANNING",             AMBER,  DKAMBER),
    ("⑧ EXECUTION",            GREEN,  DKGREEN),
    ("⑨ REVIEW",               GREEN,  DKGREEN),
    ("⑩ DELIVERY",             GREEN,  DKGREEN),
]
sw = 4.8; sh = 0.42
sx = 0.35
y = y0
for step_n, (label_t, fc, bg_c) in enumerate(steps):
    is_lock = label_t.startswith("🔒")
    bw2 = Pt(2) if is_lock else Pt(0.7)
    rect(s7, sx, y, sw, sh, fill=bg_c, border=fc, bw=bw2)
    sz = 13 if is_lock else 11
    bold = is_lock
    text(s7, label_t, sx+0.12, y+0.08, sw-0.24, 0.28,
         size=sz, bold=bold, color=fc)
    if step_n < len(steps)-1:
        text(s7, "▼", sx+sw/2-0.15, y+sh, 0.4, 0.2,
             size=8, color=ION, align=PP_ALIGN.CENTER)
    y += sh + 0.2

# Right side: explanations
rx = 5.5; rw = 7.5

# Before lock
rect(s7, rx, y0, rw, 1.62, fill=CARD, border=MUTED, bw=Pt(0.6))
text(s7, "BEFORE MISSION LOCK", rx+0.15, y0+0.08, rw-0.3, 0.28,
     size=11, bold=True, color=MUTED)
text(s7,
     "Hades can reason, ask clarifying questions, push back on vague "
     "requests, and propose a plan — but CANNOT invoke any tools or "
     "execute any commands on the system. The Partner Brain is in full control.",
     rx+0.15, y0+0.40, rw-0.3, 1.1, size=10, color=SLATE, wrap=True)

# The lock itself
rect(s7, rx, y0+1.8, rw, 1.25, fill=ION2, border=ION, bw=Pt(1.5))
text(s7, "🔒  MISSION LOCK — THE AUTHORIZATION GATE",
     rx+0.15, y0+1.92, rw-0.3, 0.32, size=12, bold=True, color=WHITE)
text(s7,
     "understanding_evaluator.py checks 4 conditions:\n"
     "  ✓  objective is populated\n"
     "  ✓  desired_outcome is populated\n"
     "  ✓  success_criteria is populated\n"
     "  ✓  mutual_understanding_reached == True\n"
     "If ALL pass → MissionStatus → AUTHORIZED_EXECUTION",
     rx+0.15, y0+2.28, rw-0.3, 0.72, size=9.5, color=WHITE, wrap=True)

# After lock
rect(s7, rx, y0+3.23, rw, 1.62, fill=CARD, border=GREEN, bw=Pt(0.6))
text(s7, "AFTER MISSION LOCK", rx+0.15, y0+3.31, rw-0.3, 0.28,
     size=11, bold=True, color=GREEN)
text(s7,
     "Executive Brain takes full control. It generates a TaskGraph, delegates "
     "to Workers, invokes Tools on the Linux system, runs the Review Engine, "
     "and delivers the result — all without the user needing to manage anything. "
     "The user retains authority over GOALS, not over process.",
     rx+0.15, y0+3.63, rw-0.3, 1.1, size=10, color=SLATE, wrap=True)

slide_num(s7, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CONVERSATION & MISSIONS
# ══════════════════════════════════════════════════════════════════════════════
s8 = slide()
top_label(s8, "HADES  —  CONVERSATION & MISSIONS")
y0 = section_title(s8, "HADES Thinks in Missions, Not Prompts.",
                   "Every piece of work is a persistent, structured object — not a disposable input.")

# Two column layout
lw = 6.1; rw2 = 6.5; lx = 0.3; rx2 = 6.73

# Left — traditional vs hades
rect(s8, lx, y0, lw, 0.32, fill=DKRED, border=RED, bw=Pt(0.6))
text(s8, "TRADITIONAL CHAT — DISPOSABLE PROMPTS", lx+0.1, y0+0.06, lw-0.2, 0.22,
     size=9, bold=True, color=RED)

trad = [
    "Prompt → Response → Context Lost",
    "New session = Start from scratch",
    "No persistent goal tracking",
    "No success criteria",
    "No outcome validation",
]
for i, t in enumerate(trad):
    yt = y0 + 0.42 + i * 0.46
    rect(s8, lx, yt, lw, 0.38, fill=DKRED, border=RGBColor(0x45,0x10,0x10), bw=Pt(0.5))
    text(s8, f"✗  {t}", lx+0.1, yt+0.08, lw-0.2, 0.24, size=10, color=SLATE)

# Right — Mission object
rect(s8, rx2, y0, rw2, 0.32, fill=DKVIOLET, border=ION, bw=Pt(0.8))
text(s8, "HADES MISSION OBJECT  —  mission.py", rx2+0.1, y0+0.06, rw2-0.2, 0.22,
     size=9, bold=True, color=GLOW)

mission_fields = [
    ("🎯  objective",           "\"Build project scaffold in /workspace\""),
    ("📋  desired_outcome",     "\"Working FastAPI app with tests\""),
    ("⚠️   constraints",        "[\"no external deps\", \"keep < 200ms\"]"),
    ("✅  success_criteria",    "\"app starts, all tests pass\""),
    ("📊  status",              "BACKGROUND_WORK"),
    ("🔒  mutual_understanding","True"),
    ("💬  conversation_history","[12 messages]"),
    ("💾  mission_id",          "\"m_3f8a9b2c\""),
]
for i, (field, val) in enumerate(mission_fields):
    ym = y0 + 0.42 + i * 0.44
    rect(s8, rx2, ym, rw2, 0.38, fill=CARD, border=RGBColor(0x30,0x30,0x55), bw=Pt(0.5))
    text(s8, field, rx2+0.1,   ym+0.08, 3.1, 0.24, size=9, color=GLOW)
    text(s8, val,   rx2+3.25,  ym+0.08, 3.15, 0.24, size=9, color=WHITE)

# State machine bar at bottom
hline(s8, 0.3, 6.62, 12.73, MUTED, Pt(0.4))
rect(s8, 0.3, 6.72, 12.73, 0.55, fill=PANEL, border=ION, bw=Pt(0.8))
text(s8, "MISSION STATE MACHINE:", 0.45, 6.81, 2.5, 0.3,
     size=9, bold=True, color=GLOW)
states_sm = [
    ("CONVERSATION", SLATE),
    ("→", ION),
    ("AUTHORIZED_EXECUTION", ION),
    ("→", ION),
    ("BACKGROUND_WORK", AMBER),
    ("→", ION),
    ("COMPLETED", GREEN),
    ("  /  NEEDS_USER", AMBER),
]
sx2 = 3.1
for st, sc in states_sm:
    sw2 = len(st) * 0.1 + 0.2
    text(s8, st, sx2, 6.81, sw2, 0.3,
         size=9, bold=(st not in ("→", "  /  NEEDS_USER")), color=sc)
    sx2 += sw2

slide_num(s8, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — REVIEW ENGINE & MEMORY
# ══════════════════════════════════════════════════════════════════════════════
s9 = slide()
top_label(s9, "HADES  —  REVIEW ENGINE & MEMORY")
y0 = section_title(s9, "AI Doesn't Get to Declare Itself Successful.",
                   "GENERATION ≠ COMPLETION  —  4-stage validation before every delivery.")

# Left: Review Engine pipeline
lx = 0.3; lw = 5.8

stages = [
    ("WORKER OUTPUT",              "LLM-generated command / artifact",     SLATE,  PANEL),
    ("STAGE 1: EXIT CODE CHECK",   "exit_code == 0 ?  (bash return value)", GREEN, DKGREEN),
    ("STAGE 2: ARTIFACT CHECK",    "Does output.txt exist on disk?",        GREEN, DKGREEN),
    ("STAGE 3: HEURISTIC CHECK",   "Non-empty, sane content validation?",   GREEN, DKGREEN),
    ("STAGE 4: SEMANTIC REVIEW",   "LLM checks output vs success_criteria", GLOW,  DKVIOLET),
]
sh2 = 0.58
y = y0
for i, (title_s, detail_s, fc, bg_c) in enumerate(stages):
    rect(s9, lx, y, lw, sh2, fill=bg_c, border=fc, bw=Pt(0.8))
    text(s9, title_s, lx+0.12, y+0.06, lw-0.24, 0.26, size=11, bold=True, color=fc)
    text(s9, detail_s, lx+0.12, y+0.32, lw-0.24, 0.2, size=9, color=SLATE)
    if i < len(stages)-1:
        text(s9, "▼", lx+lw/2-0.15, y+sh2, 0.4, 0.2,
             size=8, color=ION, align=PP_ALIGN.CENTER)
    y += sh2 + 0.2

# Pass / Retry
rect(s9, lx, y, 2.65, 0.52, fill=DKGREEN, border=GREEN, bw=Pt(1.0))
text(s9, "✓  PASS → DELIVER", lx+0.1, y+0.12, 2.45, 0.3,
     size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(s9, lx+2.95, y, 2.65, 0.52, fill=DKAMBER, border=AMBER, bw=Pt(1.0))
text(s9, "🔄  RETRY → RE-PROMPT", lx+3.05, y+0.12, 2.45, 0.3,
     size=12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Right: Memory system
rx = 6.45; rw = 6.55

rect(s9, rx, y0, rw, 0.32, fill=DKVIOLET, border=ION, bw=Pt(0.8))
text(s9, "MEMORY SYSTEM", rx+0.1, y0+0.06, rw-0.2, 0.22,
     size=10, bold=True, color=GLOW)

mem_cards = [
    ("SESSION MEMORY", GREEN,
     ["Stored in Python dict keyed by UUID",
      "Contains full conversation history",
      "MissionUnderstanding fields",
      "Active mission state",
      "Lost on server restart (in-memory)"]),
    ("MISSION HISTORY", AMBER,
     ["Appended to memory.json on completion",
      "Contains: objective, outcome, artifacts",
      "Injected into PartnerBrain system prompt",
      "Enables cross-mission awareness",
      "RAG/vector retrieval = future milestone"]),
    ("LOCAL TTS OUTPUT", GREEN,
     ["Kokoro-ONNX — local neural synthesis",
      "Audio generated server-side on Linux",
      "Encoded as Base64 in JSON response",
      "HTML5 Audio plays in browser",
      "No cloud TTS dependency"]),
]
for i, (title_m, tc, lines) in enumerate(mem_cards):
    ym = y0 + 0.42 + i * 1.6
    rect(s9, rx, ym, rw, 1.48, fill=CARD, border=tc, bw=Pt(0.7))
    text(s9, title_m, rx+0.12, ym+0.06, rw-0.24, 0.28,
         size=11, bold=True, color=tc)
    text(s9, "\n".join(lines), rx+0.12, ym+0.38, rw-0.24, 1.0,
         size=9, color=SLATE, wrap=True)

slide_num(s9, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — INNOVATION
# ══════════════════════════════════════════════════════════════════════════════
s10 = slide()
top_label(s10, "INNOVATION")
y0 = section_title(s10, "What HADES Does That No One Else Does.",
                   "A verified, 13-step end-to-end execution trace — from human speech to file on disk.")

# Mission statement in terminal style
rect(s10, 0.3, y0, 12.73, 0.48, fill=RGBColor(0x08,0x08,0x1e), border=GLOW, bw=Pt(0.8))
text(s10, ">  \"List all files in the current directory and save them to output.txt\"",
     0.5, y0+0.1, 12.3, 0.3, size=12, color=GLOW)

y0 += 0.62

# 13 steps in 2 columns
steps_e2e = [
    ("01","USER",     "Types or speaks the command",                          SLATE),
    ("02","FRONTEND", "HadesService.sendMessage() → POST /api/chat",          ION),
    ("03","API",      "main.py receives → looks up SessionState by UUID",      ION),
    ("04","PARTNER",  "IntentClassifier classifies as SMALL_TASK",             GLOW),
    ("05","PARTNER",  "MissionExtractor populates all MissionUnderstanding fields", GLOW),
    ("06","LOCK",     "UnderstandingEvaluator → MissionStatus.AUTHORIZED_EXECUTION", WHITE),
    ("07","API",      "asyncio.create_task(execution_brain) — HTTP returns immediately", ION),
    ("08","EXEC",     "ExecutionBrain._generate_plan() → 1-task TaskGraph",    AMBER),
    ("09","WORKER",   "WorkerManager → Gemini → LLM generates: ls -la > output.txt", GREEN),
    ("10","TOOL",     "TerminalSkill → asyncio.create_subprocess_shell → runs on Linux", GREEN),
    ("11","REVIEW",   "ReviewEngine: exit_code==0 ✓  · output.txt exists ✓ → PASS", GREEN),
    ("12","MEMORY",   "memory_manager.add_mission_to_history() → memory.json", SLATE),
    ("13","SSE",      "EventSource emits MISSION_COMPLETED → UI + Kokoro TTS speaks", GREEN),
]

col1 = steps_e2e[:7]
col2 = steps_e2e[7:]
sh3  = 0.44

for col_idx, col_steps in enumerate([col1, col2]):
    bx = 0.3 if col_idx == 0 else 6.72
    bw3 = 6.28
    y_s = y0
    for num, layer, action, fc in col_steps:
        rect(s10, bx, y_s, bw3, sh3, fill=CARD, border=fc, bw=Pt(0.6))
        # Number badge
        rect(s10, bx, y_s, 0.38, sh3, fill=fc, border=None)
        text(s10, num, bx, y_s+0.1, 0.38, 0.25,
             size=9, bold=True, color=BG if fc != SLATE else WHITE,
             align=PP_ALIGN.CENTER)
        # Layer badge
        rect(s10, bx+0.42, y_s+0.08, 1.1, 0.28, fill=PANEL, border=fc, bw=Pt(0.5))
        text(s10, layer, bx+0.45, y_s+0.12, 1.04, 0.2,
             size=7.5, bold=True, color=fc, align=PP_ALIGN.CENTER)
        # Action
        text(s10, action, bx+1.6, y_s+0.07, bw3-1.72, 0.3,
             size=9, color=WHITE, wrap=False)
        y_s += sh3 + 0.04

ABYSS = RGBColor(0x09, 0x09, 0x0f)

# Bottom chain
hline(s10, 0.3, 7.02, 12.73, ION, Pt(0.5))
chain_txt = "CONVERSATION  →  MISSION LOCK  →  ASYNC EXECUTION  →  VERIFICATION  →  DELIVERY"
text(s10, chain_txt, 0, 7.1, 13.33, 0.32,
     size=12, bold=True, color=GLOW, align=PP_ALIGN.CENTER)

slide_num(s10, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — NOVELTY
# ══════════════════════════════════════════════════════════════════════════════
s11 = slide()
top_label(s11, "NOVELTY")
y0 = section_title(s11, "What HADES Actually Contributes.",
                   "7 architectural novelties — each grounded in the implemented codebase.")

novelties = [
    ("01", "🔒", "MUTUAL UNDERSTANDING GATE",
     "Execution gated by deterministic state check: all mission fields + "
     "mutual_understanding_reached=True. No other system requires explicit "
     "alignment verification before tool invocation.",
     "understanding_evaluator.py", GREEN),
    ("02", "🎯", "MISSION-CENTRIC COMPUTING",
     "Prompts are disposable. Missions are persistent objects with objective, "
     "constraints, success criteria, and outcome. State persists across "
     "conversation → execution → review → memory.",
     "mission.py", GREEN),
    ("03", "🧠", "PARTNER / EXECUTIVE DUALITY",
     "Partner Brain handles conversation, never touches tools. "
     "Executive Brain handles execution, never engages the user. "
     "Prevents the classic agent contamination failure mode.",
     "partner_brain.py + execution_brain.py", GREEN),
    ("04", "⚙️", "INVISIBLE ORCHESTRATION",
     "User says a goal. Hades selects the model, routes the prompt, "
     "runs the tool, validates the result — invisibly. "
     "Intelligence becomes infrastructure.",
     "worker_manager.py + config.json", GREEN),
    ("05", "✅", "VERIFICATION-FIRST EXECUTION",
     "4-stage Review Engine: exit code → artifact check → heuristic → "
     "LLM semantic review. Worker output is not trusted. "
     "Mission criteria define success, not the AI.",
     "review_engine.py", GREEN),
    ("06", "🔄", "MODEL-AGNOSTIC AI LAYER",
     "LiteLLM wraps Gemini, Groq, OpenRouter, Ollama behind one interface. "
     "Auto-fallback. New providers via config.json. "
     "No architectural coupling to any vendor.",
     "worker_manager.py + LiteLLM", GREEN),
    ("07", "🤝", "PERSISTENT RELATIONSHIP",
     "Mission history, user identity, conversation context persist across "
     "interactions. Hades builds continuity beyond session-scoped Q&A. "
     "RAG/vector retrieval is the next memory milestone.",
     "memory_manager.py  (partial)", AMBER),
]

# 4 + 3 grid
positions_n = [
    (0.3,  1.72, 3.25, 1.6),
    (3.65, 1.72, 3.25, 1.6),
    (7.0,  1.72, 3.25, 1.6),
    (10.35,1.72, 2.88, 1.6),
    (0.3,  3.45, 4.32, 1.72),
    (4.72, 3.45, 4.32, 1.72),
    (9.14, 3.45, 4.09, 1.72),
]
for (nov, pos) in zip(novelties, positions_n):
    num, icon, title, body, ref, color = nov
    xn, yn, wn, hn = pos
    rect(s11, xn, yn, wn, hn, fill=CARD, border=color, bw=Pt(0.7))
    text(s11, f"{num} {icon}", xn+0.1, yn+0.06, 0.9, 0.26,
         size=9, bold=True, color=MUTED)
    text(s11, title, xn+0.1, yn+0.34, wn-0.2, 0.3,
         size=10, bold=True, color=color)
    text(s11, body, xn+0.1, yn+0.66, wn-0.2, hn-0.84,
         size=8, color=SLATE, wrap=True)
    text(s11, ref, xn+0.1, yn+hn-0.2, wn-0.2, 0.18,
         size=6.5, color=MUTED, italic=True)

# Bottom
hline(s11, 0.3, 5.32, 12.73, MUTED, Pt(0.5))
text(s11, "The novelty is not another smarter model.",
     0, 5.42, 13.33, 0.38, size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s11, "It is a system that organizes intelligence around human intent.",
     0, 5.82, 13.33, 0.42, size=17, bold=True, color=GLOW, align=PP_ALIGN.CENTER)

slide_num(s11, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s12 = slide()

# Top strip
rect(s12, 0, 0, 13.33, 0.08, fill=ION2, border=None)

# Central: HADES
text(s12, "H A D E S", 0, 0.55, 13.33, 1.6,
     size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

hline(s12, 1.0, 2.1, 11.33, ION, Pt(2))

# Thank You
text(s12, "THANK YOU", 0, 2.22, 13.33, 0.72,
     size=38, bold=True, color=GLOW, align=PP_ALIGN.CENTER)

hline(s12, 1.0, 2.98, 11.33, MUTED, Pt(0.5))

# Summary statement
text(s12,
     "HADES is a Linux-native AI Operating System that converts human intent "
     "into verified system-level action through conversation-first mission planning, "
     "multi-model orchestration, and 4-stage output validation.",
     1.0, 3.1, 11.33, 0.85, size=13, color=SLATE, align=PP_ALIGN.CENTER, wrap=True)

# 4 key cards
key_cards = [
    ("BUILT ON LINUX",     "bash · asyncio · Kokoro-ONNX\nfilesystem · subprocess · Ollama", GREEN),
    ("MISSION LOCKED",     "Conversation gates execution.\nNo action without alignment.", WHITE),
    ("MODEL AGNOSTIC",     "Gemini · Groq · OpenRouter\nOllama · LiteLLM fallback", GLOW),
    ("VERIFIED DELIVERY",  "4-stage Review Engine.\nGeneration ≠ Success.", AMBER),
]
cw4 = 3.1; gap4 = 0.11
x4 = (13.33 - (4*cw4 + 3*gap4)) / 2
y4 = 4.12
for title_k, body_k, color_k in key_cards:
    rect(s12, x4, y4, cw4, 1.1, fill=CARD, border=color_k, bw=Pt(0.9))
    text(s12, title_k, x4+0.1, y4+0.1,  cw4-0.2, 0.28, size=11, bold=True, color=color_k)
    text(s12, body_k,  x4+0.1, y4+0.42, cw4-0.2, 0.62, size=9,  color=SLATE, wrap=True)
    x4 += cw4 + gap4

# Vision line
hline(s12, 1.0, 5.45, 11.33, ION, Pt(1.5))
text(s12, "HUMAN SETS THE MISSION.  HADES ORCHESTRATES THE INTELLIGENCE.",
     0, 5.57, 13.33, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
hline(s12, 1.0, 6.12, 11.33, ION, Pt(1.5))

text(s12, "Today HADES proves the loop.  The next phase proves the system.",
     0, 6.25, 13.33, 0.38, size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

text(s12, "Project HADES  ·  AI Operating System Prototype  ·  Linux-Native  ·  2025",
     0, 6.75, 13.33, 0.32, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom strip
rect(s12, 0, 7.38, 13.33, 0.12, fill=ION2, border=None)

slide_num(s12, 12)


# ── SAVE ──────────────────────────────────────────────────────────────────────
out = r"d:\HACK O HADES\HackOHades.pptx"
prs.save(out)
print(f"Saved: {out}")
