import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

def add_transition(slide):
    # Add a simple fade transition
    transition = OxmlElement('p:transition')
    fade = OxmlElement('p:fade')
    transition.append(fade)
    # The p:transition element must appear after p:cSld and p:clrMapOvr, usually around index 2 or 3
    slide._element.insert(2, transition)

def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def set_bg_color(slide, hex_color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def add_text(slide, text, left, top, width, height, font_size=18, font_color="FFFFFF", bold=False, italic=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor.from_string(font_color)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    return txBox

def add_box(slide, left, top, width, height, text="", bg_color="112240", border_color="00D2FF", font_size=14, font_color="FFFFFF", bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(bg_color)
    if border_color:
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string(font_color)
        p.font.bold = bold
        p.font.name = "Segoe UI"
    return shape

def add_arrow(slide, start_x, start_y, end_x, end_y, color="8892B0"):
    # Using a line with an arrow head
    connector = slide.shapes.add_connector(
        MSO_SHAPE.DOWN_ARROW, start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = RGBColor.from_string(color)
    connector.line.width = Pt(2)
    return connector

def create_base_slide(prs):
    blank_layout = prs.slide_layouts[6] # completely blank
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, "0B0F19")
    try:
        add_transition(slide)
    except:
        pass
    return slide

def add_title(slide, text, subtitle=""):
    add_text(slide, text, Inches(0.5), Inches(0.5), Inches(12), Inches(1), font_size=36, font_color="FFFFFF", bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5), font_size=18, font_color="00D2FF")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------------------------------------------------------
# SLIDE 1 — TITLE
# ---------------------------------------------------------
s1 = create_base_slide(prs)
add_text(s1, "HADES", Inches(1), Inches(2.5), Inches(11.333), Inches(1.5), font_size=72, font_color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "AI that works with you — not just talks to you.", Inches(1), Inches(4.0), Inches(11.333), Inches(0.8), font_size=24, font_color="8892B0", align=PP_ALIGN.CENTER)
add_text(s1, "Open Innovation — Linux Based", Inches(1), Inches(4.8), Inches(11.333), Inches(0.5), font_size=16, font_color="00D2FF", align=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 2 — THE PROBLEM
# ---------------------------------------------------------
s2 = create_base_slide(prs)
add_title(s2, "The problem isn't access to AI. It's orchestration.")
add_text(s2, "Today, users have to become the orchestration layer.", Inches(0.5), Inches(1.3), Inches(12), Inches(0.5), font_size=20, font_color="8892B0")

# Flow: USER -> AI -> TOOL -> COPY/PASTE -> AI -> COMMAND -> VERIFY -> USER
items = ["USER", "AI", "TOOL", "COPY / PASTE", "ANOTHER AI", "COMMAND", "VERIFY", "USER"]
start_x = 0.5
for i, item in enumerate(items):
    add_box(s2, Inches(start_x + i*1.6), Inches(3.0), Inches(1.3), Inches(0.8), text=item, bg_color="112240", border_color="FF4444" if i in [3,5,6] else "8892B0", font_size=12, bold=True)
    if i < len(items) - 1:
        add_text(s2, "→", Inches(start_x + i*1.6 + 1.3), Inches(3.2), Inches(0.3), Inches(0.5), font_size=20, font_color="FF4444", align=PP_ALIGN.CENTER)

add_text(s2, "The Burden:", Inches(0.5), Inches(5.0), Inches(2), Inches(0.5), font_size=24, font_color="FFFFFF", bold=True)
burdens = ["Plan.", "Choose tools.", "Maintain context.", "Execute.", "Monitor.", "Verify."]
for i, b in enumerate(burdens):
    add_box(s2, Inches(2.5 + i*1.7), Inches(5.0), Inches(1.5), Inches(0.8), text=b, bg_color="1B263B", border_color=None, font_color="8892B0", font_size=14)

# ---------------------------------------------------------
# SLIDE 3 — THE SHIFT
# ---------------------------------------------------------
s3 = create_base_slide(prs)
add_title(s3, "From asking AI questions → to delegating outcomes.")

# Left Side
add_text(s3, "TRADITIONAL AI", Inches(1), Inches(2.0), Inches(5), Inches(0.5), font_size=24, font_color="8892B0", bold=True, align=PP_ALIGN.CENTER)
add_text(s3, "\"Tell me how to do it.\"", Inches(1), Inches(2.5), Inches(5), Inches(0.5), font_size=18, font_color="FFFFFF", italic=True, align=PP_ALIGN.CENTER)
add_box(s3, Inches(2.5), Inches(3.5), Inches(2), Inches(0.8), text="PROMPT", bg_color="112240")
add_text(s3, "↓", Inches(3.3), Inches(4.3), Inches(0.4), Inches(0.5), font_size=20, align=PP_ALIGN.CENTER)
add_box(s3, Inches(2.5), Inches(4.8), Inches(2), Inches(0.8), text="RESPONSE", bg_color="112240")

# Divider
add_box(s3, Inches(6.6), Inches(2.0), Inches(0.02), Inches(4.5), bg_color="8892B0", border_color=None)

# Right Side
add_text(s3, "HADES", Inches(7), Inches(2.0), Inches(5), Inches(0.5), font_size=24, font_color="00D2FF", bold=True, align=PP_ALIGN.CENTER)
add_text(s3, "\"Handle it.\"", Inches(7), Inches(2.5), Inches(5), Inches(0.5), font_size=18, font_color="FFFFFF", italic=True, align=PP_ALIGN.CENTER)
add_box(s3, Inches(8.5), Inches(3.2), Inches(2), Inches(0.6), text="GOAL", bg_color="00D2FF", font_color="0B0F19", bold=True)
add_text(s3, "↓", Inches(9.3), Inches(3.8), Inches(0.4), Inches(0.4), font_size=16, align=PP_ALIGN.CENTER)
add_box(s3, Inches(8.5), Inches(4.2), Inches(2), Inches(0.6), text="ORCHESTRATION\n&\nEXECUTION", bg_color="112240", border_color="00D2FF")
add_text(s3, "↓", Inches(9.3), Inches(4.8), Inches(0.4), Inches(0.4), font_size=16, align=PP_ALIGN.CENTER)
add_box(s3, Inches(8.0), Inches(5.2), Inches(3), Inches(0.8), text="VERIFIED RESULT", bg_color="005B4F", border_color="00E676", bold=True)

# ---------------------------------------------------------
# SLIDE 4 — WHAT IS HADES?
# ---------------------------------------------------------
s4 = create_base_slide(prs)
add_title(s4, "What is HADES?")
add_text(s4, "HADES is an AI-powered application designed to act as a working partner — understanding user goals, coordinating specialized workers and tools, executing tasks, and returning with the result.", Inches(0.5), Inches(1.3), Inches(12), Inches(1.0), font_size=20, font_color="FFFFFF")

add_box(s4, Inches(5.6), Inches(2.5), Inches(2), Inches(0.8), text="USER", bg_color="112240", border_color="8892B0")
add_text(s4, "↓", Inches(6.5), Inches(3.3), Inches(0.4), Inches(0.4), font_size=20, align=PP_ALIGN.CENTER)
add_box(s4, Inches(5.6), Inches(3.7), Inches(2), Inches(0.8), text="HADES", bg_color="112240", border_color="00D2FF", bold=True)
add_text(s4, "↓", Inches(6.5), Inches(4.5), Inches(0.4), Inches(0.4), font_size=20, align=PP_ALIGN.CENTER)
add_box(s4, Inches(4.6), Inches(4.9), Inches(4), Inches(0.8), text="WORKERS + TOOLS", bg_color="1B263B", border_color=None)
add_text(s4, "↓", Inches(6.5), Inches(5.7), Inches(0.4), Inches(0.4), font_size=20, align=PP_ALIGN.CENTER)
add_box(s4, Inches(4.6), Inches(6.1), Inches(1.8), Inches(0.8), text="LINUX", bg_color="2C1E3D", border_color="B388FF")
add_box(s4, Inches(6.8), Inches(6.1), Inches(1.8), Inches(0.8), text="RESULT", bg_color="005B4F", border_color="00E676")

# ---------------------------------------------------------
# SLIDE 5 — HOW HADES WORKS
# ---------------------------------------------------------
s5 = create_base_slide(prs)
add_title(s5, "How HADES Works", "An execution pipeline designed for outcomes.")

pipeline = ["UNDERSTAND", "PLAN", "DELEGATE", "EXECUTE", "VERIFY", "RETURN"]
for i, step in enumerate(pipeline):
    add_box(s5, Inches(1.0 + i*1.9), Inches(2.5), Inches(1.6), Inches(0.8), text=step, bg_color="112240", border_color="00D2FF", font_size=12, bold=True)
    if i < len(pipeline) - 1:
        add_text(s5, "→", Inches(1.0 + i*1.9 + 1.6), Inches(2.65), Inches(0.3), Inches(0.5), font_size=20, align=PP_ALIGN.CENTER)

add_text(s5, "User Visibility:", Inches(1.0), Inches(4.0), Inches(3), Inches(0.5), font_size=24, font_color="FFFFFF", bold=True)
messages = [
    "\"Got it. I'm on it.\"",
    "\"I've handled the research.\"",
    "\"I'm working through the files.\"",
    "\"I'll ping you when it's done.\""
]
for i, msg in enumerate(messages):
    add_box(s5, Inches(1.0), Inches(4.8 + i*0.6), Inches(4.0), Inches(0.4), text=msg, bg_color="0B0F19", border_color="8892B0", font_size=14, font_color="8892B0", align=PP_ALIGN.LEFT)

add_text(s5, "The user remains in control but does not have to micromanage execution. No chain-of-thought walls of text.", Inches(6.0), Inches(5.0), Inches(6.0), Inches(1.0), font_size=20, font_color="00D2FF", italic=True)

# ---------------------------------------------------------
# SLIDE 6 — PARTNER BRAIN + EXECUTION BRAIN
# ---------------------------------------------------------
s6 = create_base_slide(prs)
add_title(s6, "Internal Architecture", "Two distinct systems for interaction and action.")

# Partner Brain
pb = add_box(s6, Inches(1.5), Inches(2.5), Inches(4), Inches(3.5), text="PARTNER BRAIN\n\n• Conversation\n• Context\n• Intent\n• Clarification\n• User relationship", bg_color="112240", border_color="8892B0", font_size=18, align=PP_ALIGN.LEFT)
# Execution Brain
eb = add_box(s6, Inches(7.5), Inches(2.5), Inches(4), Inches(3.5), text="EXECUTION BRAIN\n\n• Mission\n• Planning\n• Orchestration\n• Delegation\n• Execution\n• Verification", bg_color="112240", border_color="00D2FF", font_size=18, align=PP_ALIGN.LEFT)

add_text(s6, "→", Inches(5.8), Inches(4.0), Inches(1.4), Inches(0.5), font_size=40, font_color="00D2FF", align=PP_ALIGN.CENTER)
add_text(s6, "Hands off intent", Inches(5.8), Inches(4.8), Inches(1.4), Inches(0.5), font_size=12, font_color="8892B0", align=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 7 — WORKERS
# ---------------------------------------------------------
s7 = create_base_slide(prs)
add_title(s7, "HADES doesn't have to do everything itself.")

add_box(s7, Inches(5.6), Inches(1.5), Inches(2), Inches(0.8), text="HADES", bg_color="112240", border_color="00D2FF", bold=True)

# Lines
add_box(s7, Inches(6.5), Inches(2.3), Inches(0.02), Inches(0.5), bg_color="8892B0", border_color=None)
add_box(s7, Inches(2.5), Inches(2.8), Inches(8.0), Inches(0.02), bg_color="8892B0", border_color=None)
add_box(s7, Inches(2.5), Inches(2.8), Inches(0.02), Inches(0.5), bg_color="8892B0", border_color=None)
add_box(s7, Inches(6.5), Inches(2.8), Inches(0.02), Inches(0.5), bg_color="8892B0", border_color=None)
add_box(s7, Inches(10.5), Inches(2.8), Inches(0.02), Inches(0.5), bg_color="8892B0", border_color=None)

# Workers
add_box(s7, Inches(1.5), Inches(3.3), Inches(2.0), Inches(0.8), text="Research Worker", bg_color="1B263B", font_color="00D2FF", bold=True)
add_box(s7, Inches(5.5), Inches(3.3), Inches(2.0), Inches(0.8), text="Computer Worker", bg_color="1B263B", font_color="00D2FF", bold=True)
add_box(s7, Inches(9.5), Inches(3.3), Inches(2.0), Inches(0.8), text="Creation Worker", bg_color="1B263B", font_color="00D2FF", bold=True)

# Worker outputs
add_box(s7, Inches(1.5), Inches(4.5), Inches(2.0), Inches(0.6), text="Web / Data", bg_color="0B0F19", border_color="8892B0")
add_box(s7, Inches(5.5), Inches(4.5), Inches(2.0), Inches(0.6), text="Linux Shell", bg_color="0B0F19", border_color="8892B0")
add_box(s7, Inches(9.5), Inches(4.5), Inches(2.0), Inches(0.6), text="Files / Docs", bg_color="0B0F19", border_color="8892B0")

add_text(s7, "Workers are specialized execution units/models.\nHADES decides which worker is appropriate.\nThe user does not need to manually select the worker.", Inches(1.0), Inches(6.0), Inches(11.333), Inches(1.0), font_size=20, font_color="FFFFFF", align=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 8 — LINUX
# ---------------------------------------------------------
s8 = create_base_slide(prs)
add_title(s8, "Why Linux?", "Operating at the Linux application layer.")

add_box(s8, Inches(1.0), Inches(2.0), Inches(2.5), Inches(1.0), text="HADES", bg_color="112240", border_color="00D2FF", font_size=20, bold=True)
for i in range(5):
    add_text(s8, "↓", Inches(2.1), Inches(3.1 + i*0.4), Inches(0.4), Inches(0.4), font_size=16, font_color="8892B0")

add_box(s8, Inches(1.0), Inches(5.5), Inches(2.5), Inches(1.0), text="LINUX\nCapabilities", bg_color="2C1E3D", border_color="B388FF", font_size=18, bold=True)

caps = ["Terminal", "Filesystem", "Processes", "Browser", "Research", "Automation", "Applications"]
for i, cap in enumerate(caps):
    add_box(s8, Inches(4.5 + (i%3)*2.8), Inches(2.5 + (i//3)*1.2), Inches(2.5), Inches(0.8), text=cap, bg_color="1B263B", font_size=16)

add_text(s8, "HADES does not replace Linux or modify the kernel.\nHADES is an AI-powered application that uses the Linux environment to perform real work.", Inches(4.5), Inches(5.5), Inches(8.0), Inches(1.0), font_size=20, font_color="00D2FF", bold=True)

# ---------------------------------------------------------
# SLIDE 9 — NOVELTY
# ---------------------------------------------------------
s9 = create_base_slide(prs)
add_title(s9, "What makes HADES different?")

cards = [
    ("PARTNER MODEL", "HADES owns the workflow instead of making the user orchestrate it."),
    ("WORKER ARCHITECTURE", "Specialized workers can handle different jobs dynamically."),
    ("OUTCOME-FIRST", "Users describe what needs to happen rather than micro-managing every step."),
    ("LINUX EXECUTION", "AI can interact with the real computing environment directly.")
]

for i, (title, desc) in enumerate(cards):
    left = 0.5 + (i%2)*6.2
    top = 2.0 + (i//2)*2.5
    add_box(s9, Inches(left), Inches(top), Inches(5.8), Inches(2.0), text=f"{title}\n\n{desc}", bg_color="112240", border_color="00D2FF" if i==0 else "8892B0", font_size=16, align=PP_ALIGN.LEFT)

# ---------------------------------------------------------
# SLIDE 10 — REAL-WORLD USE CASE
# ---------------------------------------------------------
s10 = create_base_slide(prs)
add_title(s10, "Real-World Use Case", "Intelligent Workspace Cleanup")

# User prompt
add_box(s10, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.8), text="USER: \"Inspect my workspace, find what's taking up space, and tell me what is safe to clean.\"", bg_color="1B263B", border_color="00D2FF", font_size=18, bold=True, align=PP_ALIGN.LEFT)

# HADES execution
add_box(s10, Inches(0.5), Inches(2.8), Inches(5.0), Inches(3.5), text="HADES Execution:\n\n✓ Understands goal\n✓ Inspects workspace\n✓ Analyzes files\n✓ Categorizes caches/logs/builds\n✓ Determines safe candidates\n✓ Verifies findings\n✓ Returns report", bg_color="112240", font_size=16, align=PP_ALIGN.LEFT)

# Result
add_box(s10, Inches(6.0), Inches(3.5), Inches(6.8), Inches(2.0), text="RESULT:\n\"Here's what I found. These files are safe candidates. These require caution. Nothing was deleted without authorization.\"", bg_color="005B4F", border_color="00E676", font_size=18, align=PP_ALIGN.LEFT)

# ---------------------------------------------------------
# SLIDE 11 — INNOVATION
# ---------------------------------------------------------
s11 = create_base_slide(prs)
add_title(s11, "Innovation: Changing the Relationship")

# Current AI
add_text(s11, "CURRENT AI", Inches(2.0), Inches(1.5), Inches(3), Inches(0.5), font_size=20, font_color="8892B0", bold=True, align=PP_ALIGN.CENTER)
add_box(s11, Inches(2.0), Inches(2.2), Inches(3), Inches(0.8), text="Conversation", bg_color="112240")
add_text(s11, "↓", Inches(3.3), Inches(3.1), Inches(0.4), Inches(0.5), font_size=24)
add_box(s11, Inches(2.0), Inches(3.7), Inches(3), Inches(0.8), text="Answer", bg_color="112240")

# HADES
add_text(s11, "HADES", Inches(8.0), Inches(1.5), Inches(3), Inches(0.5), font_size=20, font_color="00D2FF", bold=True, align=PP_ALIGN.CENTER)
steps = ["Intent", "Planning", "Workers", "Tools", "Linux", "Verification", "Outcome"]
for i, step in enumerate(steps):
    add_box(s11, Inches(7.5), Inches(2.1 + i*0.7), Inches(4), Inches(0.5), text=step, bg_color="1B263B" if i<6 else "005B4F", border_color="00D2FF" if i<6 else "00E676", font_size=14)
    if i < len(steps)-1:
        add_text(s11, "↓", Inches(9.3), Inches(2.65 + i*0.7), Inches(0.4), Inches(0.2), font_size=12)

# ---------------------------------------------------------
# SLIDE 12 — CURRENT IMPLEMENTATION
# ---------------------------------------------------------
s12 = create_base_slide(prs)
add_title(s12, "Implementation Status", "Honesty in execution.")

add_box(s12, Inches(1.0), Inches(2.0), Inches(5.0), Inches(4.0), text="CURRENTLY IMPLEMENTED\n\n• Core Partner Brain loop\n• Execution Brain orchestrator\n• File & Shell Worker integration\n• Autonomous tool usage\n• LLM API integration (providers)\n• Safe execution verification", bg_color="112240", border_color="00E676", font_size=18, align=PP_ALIGN.LEFT)

add_box(s12, Inches(7.333), Inches(2.0), Inches(5.0), Inches(4.0), text="PLANNED / NEXT\n\n• Hardened sandbox execution\n• Advanced visual/browser workers\n• Self-healing failure recovery\n• Voice and multimodal intent\n• Long-term episodic memory\n• Custom evaluation datasets", bg_color="112240", border_color="8892B0", font_size=18, font_color="8892B0", align=PP_ALIGN.LEFT)

# ---------------------------------------------------------
# SLIDE 13 — LIVE DEMO
# ---------------------------------------------------------
s13 = create_base_slide(prs)
add_text(s13, "One mission. End to end.", Inches(1.0), Inches(1.5), Inches(11.333), Inches(1.0), font_size=48, font_color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)

flow = ["USER", "GOAL", "HADES", "MISSION", "WORKERS", "LINUX", "EXECUTION", "VERIFICATION", "RESULT"]
start_x = 0.5
for i, item in enumerate(flow):
    add_box(s13, Inches(start_x + i*1.4), Inches(4.0), Inches(1.2), Inches(0.8), text=item, bg_color="112240", border_color="00D2FF", font_size=10, bold=True)
    if i < len(flow) - 1:
        add_text(s13, "→", Inches(start_x + i*1.4 + 1.1), Inches(4.2), Inches(0.3), Inches(0.3), font_size=16, font_color="8892B0", align=PP_ALIGN.CENTER)

add_text(s13, "Transitioning to Live Demo...", Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.8), font_size=20, font_color="00D2FF", italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 14 — CLOSING
# ---------------------------------------------------------
s14 = create_base_slide(prs)
add_text(s14, "Tell HADES what needs to happen.\nLet it handle the work.", Inches(1.0), Inches(2.5), Inches(11.333), Inches(1.5), font_size=40, font_color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)

# Final summary box
add_box(s14, Inches(2.666), Inches(5.0), Inches(8.0), Inches(1.0), text="AI Partner  +  Workers  +  Linux  +  Execution", bg_color="112240", border_color="00D2FF", font_size=24, bold=True)

prs.save('HADES_Hackathon_Presentation.pptx')
print("Generated HADES_Hackathon_Presentation.pptx")
