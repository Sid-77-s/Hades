from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── COLOUR PALETTE ────────────────────────────────────────────────────────────
VOID    = RGBColor(0x09, 0x09, 0x0f)   # near-black bg
ABYSS   = RGBColor(0x0f, 0x0f, 0x1a)   # slide bg
PANEL   = RGBColor(0x14, 0x14, 0x2a)   # card bg
ION     = RGBColor(0x7c, 0x3a, 0xed)   # violet accent
SIGNAL  = RGBColor(0x6d, 0x28, 0xd9)   # darker violet
GLOW    = RGBColor(0xa7, 0x8b, 0xfa)   # lavender highlights
WHITE   = RGBColor(0xff, 0xff, 0xff)
ONLINE  = RGBColor(0x10, 0xb9, 0x81)   # green – IMPLEMENTED
AMBER   = RGBColor(0xf5, 0x9e, 0x0b)   # amber – PARTIAL
DANGER  = RGBColor(0xef, 0x44, 0x44)   # red – MISSING
MUTED   = RGBColor(0x6b, 0x72, 0x80)
SUBTLE  = RGBColor(0x94, 0xa3, 0xb8)
TEXT    = RGBColor(0xe2, 0xe8, 0xf0)
DARK_VIOLET = RGBColor(0x1e, 0x10, 0x40)  # deep panel bg

# Slide dimensions (16:9 widescreen)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank layout


# ── HELPER FUNCTIONS ──────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    # Dark background fill
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = ABYSS
    return slide


def box(slide, x, y, w, h, fill=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = border_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def label(slide, text, x, y, w, h, bg, fg, size=9):
    """Coloured badge label."""
    b = box(slide, x, y, w, h, fill=bg, border_color=fg, border_width=Pt(0.75))
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = False
    tf = t.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Calibri"
    return b


def pill_row(slide, items, y, colors):
    """Horizontal row of pill-shaped items."""
    n = len(items)
    gap = 0.15
    total_w = 13.0
    w_each = (total_w - gap * (n - 1)) / n
    x = 0.17
    for i, item in enumerate(items):
        c = colors[i % len(colors)]
        box(slide, x, y, w_each, 0.4, fill=PANEL, border_color=c)
        txt(slide, item, x + 0.05, y + 0.05, w_each - 0.1, 0.3,
            size=11, color=c, align=PP_ALIGN.CENTER, bold=True)
        x += w_each + gap


def hline(slide, x, y, w, color=ION, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def section_header(slide, title, sub=None, y_title=0.25):
    txt(slide, title, 0.4, y_title, 12.5, 0.7,
        size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if sub:
        txt(slide, sub, 0.4, y_title + 0.65, 12.5, 0.4,
            size=18, bold=False, color=GLOW, align=PP_ALIGN.LEFT)
    hline(slide, 0.4, y_title + (1.1 if sub else 0.75), 12.5, color=ION)


def arrow_down(slide, x, y):
    txt(slide, "▼", x, y, 0.4, 0.3, size=10, color=ION, align=PP_ALIGN.CENTER)


def status_badge(slide, status, x, y):
    if status == "IMPLEMENTED":
        label(slide, "✓  IMPLEMENTED", x, y, 1.6, 0.28,
              bg=RGBColor(0x05, 0x2e, 0x1a), fg=ONLINE, size=8)
    elif status == "PARTIAL":
        label(slide, "⚡  PARTIAL", x, y, 1.3, 0.28,
              bg=RGBColor(0x2e, 0x1f, 0x05), fg=AMBER, size=8)
    elif status == "FUTURE":
        label(slide, "◆  FUTURE", x, y, 1.2, 0.28,
              bg=RGBColor(0x12, 0x0a, 0x2e), fg=GLOW, size=8)
    elif status == "MISSING":
        label(slide, "✗  MISSING", x, y, 1.2, 0.28,
              bg=RGBColor(0x2e, 0x08, 0x08), fg=DANGER, size=8)


def card(slide, x, y, w, h, title, body_lines, title_color=GLOW,
         bg=PANEL, border=ION, title_size=13, body_size=10):
    box(slide, x, y, w, h, fill=bg, border_color=border)
    txt(slide, title, x + 0.12, y + 0.08, w - 0.24, 0.28,
        size=title_size, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txt(slide, body, x + 0.12, y + 0.36, w - 0.24, h - 0.46,
        size=body_size, color=TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()

# Badge top
label(s1, "HACKATHON BUILD  ·  AI SYSTEMS ENGINEERING  ·  2025",
      3.2, 0.22, 6.9, 0.32, bg=DARK_VIOLET, fg=GLOW, size=9)

# Giant HADES
txt(s1, "HADES", 0.5, 0.8, 12.3, 2.2,
    size=120, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Violet accent line under title
hline(s1, 1.5, 2.9, 10.3, color=ION, thickness=Pt(2))

# Subtitle
txt(s1, "FROM AI ASSISTANT  →  AI OPERATING SYSTEM",
    0.5, 3.0, 12.3, 0.55,
    size=17, bold=False, color=SUBTLE, align=PP_ALIGN.CENTER)

# Pipeline row
pipeline = [
    ("HUMAN", SUBTLE),
    ("→→", ION),
    ("HADES", WHITE),
    ("→→", ION),
    ("MODELS + TOOLS + SYSTEM", GLOW),
    ("→→", ION),
    ("REAL RESULT", ONLINE),
]
x_pos = 0.4
for text, color in pipeline:
    w_node = 1.9 if text not in ("→→",) else 0.45
    if text == "HADES":
        box(s1, x_pos, 3.75, w_node, 0.52, fill=ION, border_color=SIGNAL)
        txt(s1, text, x_pos + 0.05, 3.82, w_node - 0.1, 0.38,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    elif text == "→→":
        txt(s1, text, x_pos, 3.8, w_node, 0.38,
            size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    else:
        box(s1, x_pos, 3.75, w_node, 0.52, fill=PANEL, border_color=color)
        txt(s1, text, x_pos + 0.05, 3.82, w_node - 0.1, 0.38,
            size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    x_pos += w_node + 0.05

# Bottom statement
hline(s1, 1.5, 4.55, 10.3, color=MUTED, thickness=Pt(0.5))
txt(s1,
    "The user should manage the goal — not the intelligence required to achieve it.",
    0.5, 4.65, 12.3, 0.45,
    size=13, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)

# Slide number
txt(s1, "01 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
section_header(s2, "THE AI PROBLEM IS NO LONGER CAPABILITY.", "IT IS ORCHESTRATION.")

# Left — fragmented tools
txt(s2, "TODAY'S REALITY", 0.4, 1.55, 5.5, 0.35, size=10, bold=True, color=MUTED)

tools = [
    ("🔍  GEMINI",    "Research",    0.4, 2.0),
    ("💻  CLAUDE",    "Coding",      2.7, 2.0),
    ("✍️   CHATGPT",  "Writing",     5.0, 2.0),
    ("📝  NOTION",    "Notes",       0.4, 2.85),
    ("⚡  ZAPIER",    "Automation",  2.7, 2.85),
    ("📅  CALENDAR",  "Scheduling",  5.0, 2.85),
    ("📧  GMAIL",     "Email",       1.6, 3.7),
    ("🔧  OTHER",     "Tools",       3.8, 3.7),
]
for name, task, tx, ty in tools:
    box(s2, tx, ty, 2.1, 0.62, fill=PANEL, border_color=MUTED)
    txt(s2, name, tx + 0.1, ty + 0.04, 1.5, 0.3, size=11, bold=True, color=TEXT)
    txt(s2, task, tx + 0.1, ty + 0.32, 1.5, 0.25, size=9, color=MUTED)

# "YOU" in the center of chaos
box(s2, 2.4, 2.85, 1.0, 0.5, fill=RGBColor(0x2e,0x08,0x08), border_color=DANGER)
txt(s2, "YOU", 2.4, 2.9, 1.0, 0.4, size=14, bold=True, color=DANGER, align=PP_ALIGN.CENTER)
txt(s2, "The Orchestration Layer", 0.3, 4.45, 5.5, 0.3, size=9, italic=True, color=MUTED)

# Divider
box(s2, 6.3, 1.5, 0.01, 5.5, fill=ION, border_color=None)

# Right — the hidden tax
txt(s2, "THE HIDDEN TAX", 6.55, 1.55, 6.5, 0.35, size=10, bold=True, color=MUTED)

problems = [
    ("Context switching between AI tools"),
    ("Repeating yourself to every model"),
    ("Manually selecting the right AI for each task"),
    ("Managing workflow sequencing yourself"),
    ("Context lost between every session"),
    ("Decision fatigue from tool sprawl"),
]
for i, prob in enumerate(problems):
    yp = 2.0 + i * 0.62
    box(s2, 6.55, yp, 6.5, 0.52, fill=RGBColor(0x2e,0x08,0x08), border_color=DANGER)
    txt(s2, f"✗  {prob}", 6.7, yp + 0.1, 6.25, 0.35, size=11, color=TEXT)

# Bottom statement
hline(s2, 0.4, 6.1, 12.5, color=MUTED, thickness=Pt(0.5))
txt(s2, "Today's AI gives humans MORE capabilities.  Hades removes the burden of COORDINATING them.",
    0.4, 6.2, 12.5, 0.45, size=13, color=GLOW, align=PP_ALIGN.CENTER, bold=True)

txt(s2, "02 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — COMPARISON MATRIX
# ══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
section_header(s3, "WHERE EXISTING AI SYSTEMS STOP.", "CAPABILITY GAP MATRIX")

rows = [
    "Conversation",
    "Persistent Identity",
    "Long-Term Memory",
    "Goal Understanding",
    "Human Alignment Gate",
    "Multi-Model Orchestration",
    "Background Execution",
    "Output Validation",
    "Human Goal Authority",
    "Continuous Relationship",
]
cols = ["TRADITIONAL AI", "AI AGENT", "HADES"]
col_vals = [
    # Traditional AI
    ["✓", "✗", "✗", "~", "✗", "✗", "✗", "✗", "✗", "✗"],
    # AI Agent
    ["~", "✗", "~", "~", "✗", "~", "✓", "✗", "✗", "✗"],
    # HADES
    ["✓", "✓", "~", "✓", "✓", "✓", "✓", "✓", "✓", "~"],
]

# Header row
col_x = [2.2, 5.4, 8.6]
col_w = [3.0, 3.0, 3.6]
for i, col in enumerate(cols):
    bg = ION if i == 2 else PANEL
    box(s3, col_x[i], 1.55, col_w[i], 0.48, fill=bg, border_color=ION)
    txt(s3, col, col_x[i] + 0.05, 1.6, col_w[i] - 0.1, 0.38,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Row headers + cells
row_h = 0.42
for r, row_name in enumerate(rows):
    y = 2.08 + r * row_h
    # row label
    row_bg = PANEL if r % 2 == 0 else RGBColor(0x10, 0x10, 0x22)
    box(s3, 0.15, y, 2.0, row_h - 0.02, fill=row_bg, border_color=RGBColor(0x30,0x30,0x50))
    txt(s3, row_name, 0.22, y + 0.05, 1.85, row_h - 0.1, size=10, color=TEXT)

    for i in range(3):
        val = col_vals[i][r]
        if val == "✓":
            fg, bg_c = ONLINE, RGBColor(0x05, 0x20, 0x12)
        elif val == "~":
            fg, bg_c = AMBER, RGBColor(0x20, 0x15, 0x03)
        else:
            fg, bg_c = DANGER, RGBColor(0x20, 0x05, 0x05)
        hades_border = ION if i == 2 else RGBColor(0x30,0x30,0x50)
        box(s3, col_x[i], y, col_w[i], row_h - 0.02, fill=bg_c, border_color=hades_border)
        txt(s3, val, col_x[i], y + 0.04, col_w[i], row_h - 0.1,
            size=16, bold=True, color=fg, align=PP_ALIGN.CENTER)

# Bottom gap box
box(s3, 0.15, 6.35, 12.95, 0.68, fill=DARK_VIOLET, border_color=ION)
txt(s3, "RESEARCH GAP:  Persistent understanding + mutual alignment + intelligence orchestration + verified execution",
    0.3, 6.46, 12.7, 0.45, size=11, bold=True, color=GLOW, align=PP_ALIGN.CENTER)

txt(s3, "03 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 04 — MISSION LOCK
# ══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
section_header(s4, "HADES DOESN'T START WITH A TASK.", "IT STARTS WITH UNDERSTANDING.")

steps = [
    ("① USER INTENT",          "INPUT",         SUBTLE,  PANEL),
    ("② CONVERSATION",          "IMPLEMENTED ✓", ONLINE,  RGBColor(0x05,0x20,0x12)),
    ("③ CLARIFICATION",         "IMPLEMENTED ✓", ONLINE,  RGBColor(0x05,0x20,0x12)),
    ("④ PUSHBACK / ALIGN",      "IMPLEMENTED ✓", ONLINE,  RGBColor(0x05,0x20,0x12)),
    ("⑤ MUTUAL UNDERSTANDING",  "IMPLEMENTED ✓", ONLINE,  RGBColor(0x05,0x20,0x12)),
]

lock_step = ("🔒  MISSION LOCK",  "IMPLEMENTED ✓", WHITE, ION)

after_steps = [
    ("⑦ PLANNING",    "PARTIAL ⚡",    AMBER, RGBColor(0x20,0x15,0x03)),
    ("⑧ EXECUTION",   "IMPLEMENTED ✓", ONLINE, RGBColor(0x05,0x20,0x12)),
    ("⑨ REVIEW",      "IMPLEMENTED ✓", ONLINE, RGBColor(0x05,0x20,0x12)),
    ("⑩ DELIVERY",    "IMPLEMENTED ✓", ONLINE, RGBColor(0x05,0x20,0x12)),
]

step_w = 4.5
step_h = 0.38
x_start = 4.3
y = 1.55

for name, badge, badge_color, bg_color in steps:
    box(s4, x_start, y, step_w, step_h, fill=bg_color, border_color=badge_color)
    txt(s4, name, x_start + 0.12, y + 0.06, 3.0, 0.28, size=11, bold=True, color=badge_color)
    label(s4, badge, x_start + 3.15, y + 0.05, 1.3, 0.27,
          bg=RGBColor(0x05,0x20,0x12) if badge_color==ONLINE else PANEL, fg=badge_color, size=7)
    txt(s4, "▼", x_start + 2.1, y + step_h, 0.4, 0.28, size=9, color=ION, align=PP_ALIGN.CENTER)
    y += step_h + 0.28

# MISSION LOCK — big highlighted bar
y_lock = y
box(s4, x_start - 0.2, y_lock, step_w + 0.4, 0.58, fill=ION, border_color=SIGNAL)
txt(s4, "🔒  MISSION LOCK", x_start, y_lock + 0.1, step_w, 0.4,
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s4, "NO EXECUTION WITHOUT MUTUAL UNDERSTANDING",
    x_start - 0.2, y_lock + 0.62, step_w + 0.4, 0.25,
    size=8, color=GLOW, align=PP_ALIGN.CENTER, italic=True)
y = y_lock + 0.92

for name, badge, badge_color, bg_color in after_steps:
    txt(s4, "▼", x_start + 2.1, y, 0.4, 0.28, size=9, color=ION, align=PP_ALIGN.CENTER)
    y += 0.28
    box(s4, x_start, y, step_w, step_h, fill=bg_color, border_color=badge_color)
    txt(s4, name, x_start + 0.12, y + 0.06, 3.0, 0.28, size=11, bold=True, color=badge_color)
    label(s4, badge, x_start + 3.15, y + 0.05, 1.3, 0.27,
          bg=bg_color, fg=badge_color, size=7)
    y += step_h

# Annotations left / right
txt(s4, "BEFORE LOCK:\nHades reasons, clarifies,\npushes back — but\nCANNOT execute tools.",
    0.2, 2.4, 3.8, 1.4, size=10, color=MUTED, italic=True)

txt(s4, "AFTER LOCK:\nExecution Brain takes\ncontrol. User controls\nthe GOAL, not the process.",
    9.2, 2.4, 3.8, 1.4, size=10, color=MUTED, italic=True)

hline(s4, 0.4, 7.0, 12.5, color=MUTED, thickness=Pt(0.5))
txt(s4, "Conversation is the authorization layer between intent and action.  |  understanding_evaluator.py · mission.py",
    0.4, 7.08, 12.5, 0.3, size=8, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

txt(s4, "04 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 05 — ARCHITECTURE LAYERS
# ══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
section_header(s5, "ONE PARTNER. MULTIPLE BRAINS.", "MANY CAPABILITIES.")

layers = [
    # (label, sublabel, badge_text, badge_color, bg, border, note)
    ("USER",              "Goal  ·  Context  ·  Decisions",
     "HUMAN IN CONTROL", SUBTLE, RGBColor(0x12,0x12,0x20), MUTED, None),
    ("PARTNER BRAIN",     "Intent Classification  |  Mission Extraction  |  Conversational Alignment  |  Session Memory",
     "IMPLEMENTED ✓", ONLINE, RGBColor(0x05,0x22,0x18), ONLINE,
     "partner_brain.py · intent_classifier.py · understanding_evaluator.py"),
    ("🔒  MISSION LOCK",  "MUTUAL UNDERSTANDING GATE",
     "IMPLEMENTED ✓", WHITE, ION, SIGNAL, None),
    ("EXECUTIVE BRAIN",   "Task Graph  |  Worker Delegation  |  Recovery & Retry",
     "PARTIAL ⚡", AMBER, RGBColor(0x20,0x15,0x03), AMBER,
     "execution_brain.py  —  currently generates 1 task; dynamic decomposition = FUTURE"),
    ("WORKERS",           "Terminal  |  Filesystem  |  Process Manager  |  Browser (partial)",
     "PARTIAL ⚡", AMBER, RGBColor(0x20,0x18,0x03), AMBER,
     "worker_manager.py · skills/  —  workers = model endpoints today; autonomous agents = FUTURE"),
    ("MODEL PROVIDERS",   "Gemini (Google)  |  Llama 3.3 (Groq)  |  OpenRouter  |  Ollama (Local)",
     "IMPLEMENTED ✓", ONLINE, RGBColor(0x04,0x1e,0x12), ONLINE,
     "config.json · LiteLLM"),
    ("REVIEW ENGINE",     "Exit Code Validation  |  Artifact Check  |  Semantic LLM Review",
     "IMPLEMENTED ✓", ONLINE, RGBColor(0x04,0x1a,0x20), RGBColor(0x10,0x80,0x80),
     "review_engine.py"),
    ("MEMORY",            "Session (in-memory)  |  Mission History (JSON)",
     "PARTIAL ⚡", AMBER, RGBColor(0x20,0x15,0x03), AMBER,
     "memory_manager.py  —  Vector/RAG = FUTURE"),
]

lh = 0.55
y = 1.55
for lbl, sub, badge_text, badge_color, bg_c, border_c, note in layers:
    box(s5, 0.15, y, 12.9, lh, fill=bg_c, border_color=border_c)
    txt(s5, lbl, 0.3, y + 0.05, 4.0, 0.28, size=12, bold=True, color=badge_color)
    txt(s5, sub, 0.3, y + 0.3, 8.5, 0.2, size=8, color=SUBTLE)
    if note:
        txt(s5, note, 0.3, y + 0.38, 8.5, 0.18, size=7, color=MUTED, italic=True)
    # Badge right side
    label(s5, badge_text, 11.3, y + 0.12, 1.7, 0.28,
          bg=bg_c, fg=badge_color, size=8)
    y += lh + 0.04

txt(s5, "05 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — MISSIONS NOT PROMPTS
# ══════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
section_header(s6, "HADES THINKS IN MISSIONS.", "NOT PROMPTS. NOT SESSIONS. NOT TABS.")

# Left — traditional loop
txt(s6, "TRADITIONAL AI", 0.2, 1.55, 5.5, 0.3, size=10, bold=True, color=MUTED)
trad_items = [
    ("USER sends Prompt", MUTED),
    ("▼", ION),
    ("AI returns Response", MUTED),
    ("▼", ION),
    ("Context lost. Start over.", DANGER),
]
y_t = 2.0
for item, color in trad_items:
    align = PP_ALIGN.CENTER if item == "▼" else PP_ALIGN.LEFT
    sz = 9 if item == "▼" else 11
    box_c = PANEL if item not in ("▼",) else None
    if item not in ("▼",):
        box(s6, 0.2, y_t, 5.5, 0.45, fill=PANEL, border_color=MUTED)
    txt(s6, item, 0.3, y_t + 0.08, 5.3, 0.3, size=sz, color=color, align=align)
    y_t += 0.58 if item != "▼" else 0.32

txt(s6, "Every message starts fresh.", 0.2, 4.15, 5.5, 0.3, size=9, italic=True, color=MUTED)

for stat in ["Context: None", "State: Disposable", "Memory: Zero"]:
    label(s6, stat, 0.2, y_t, 2.2, 0.28, bg=RGBColor(0x2e,0x08,0x08), fg=DANGER, size=8)
    y_t += 0.36

# Divider
box(s6, 6.0, 1.5, 0.015, 5.8, fill=ION, border_color=None)

# Right — Mission Object
txt(s6, "HADES MISSION OBJECT", 6.2, 1.55, 6.9, 0.3, size=10, bold=True, color=GLOW)
box(s6, 6.2, 1.9, 6.9, 3.9, fill=RGBColor(0x10,0x0c,0x28), border_color=ION)

mission_fields = [
    ("🎯  objective",            "\"Build a project scaffold...\""),
    ("📋  desired_outcome",      "\"Working directory with...\""),
    ("⚠️   constraints",         "[\"No external deps\", ...]"),
    ("✅  success_criteria",      "\"Directory exists, tests pass\""),
    ("📊  status",               "BACKGROUND_WORK"),
    ("🔒  mutual_understanding", "true"),
    ("💬  conversation_history", "[14 messages]"),
    ("💾  mission_id",           "\"m_8f3a...\""),
    ("🏆  output_artifacts",     "[\"scaffold.zip\", ...]"),
]
y_m = 2.0
for field, val in mission_fields:
    txt(s6, field, 6.35, y_m, 3.2, 0.32, size=9, color=GLOW)
    txt(s6, val,   9.6,  y_m, 3.4, 0.32, size=9, color=TEXT)
    y_m += 0.38

txt(s6, "src/models/mission.py", 6.2, 5.82, 6.9, 0.25, size=7, color=MUTED, italic=True)

# State machine bar
box(s6, 0.15, 6.0, 12.9, 0.52, fill=PANEL, border_color=ION)
states = [
    ("CONVERSATION", MUTED),
    ("→", ION),
    ("AUTHORIZED_EXECUTION", ION),
    ("→", ION),
    ("BACKGROUND_WORK", AMBER),
    ("→", ION),
    ("COMPLETED", ONLINE),
    ("/ NEEDS_USER", AMBER),
]
sx = 0.3
for state, color in states:
    w_s = len(state) * 0.11 + 0.3
    txt(s6, state, sx, 6.12, w_s, 0.3, size=9, bold=(state not in ("→", "/")), color=color)
    sx += w_s

txt(s6, "06 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 07 — WORKER MANAGER / MODEL ROUTING
# ══════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
section_header(s7, "HADES MANAGES AI.", "THE USER DOESN'T.")

# Center box — WorkerManager
box(s7, 4.9, 2.8, 3.5, 1.2, fill=DARK_VIOLET, border_color=ION)
txt(s7, "WORKER MANAGER", 4.9, 2.88, 3.5, 0.45, size=14, bold=True, color=GLOW, align=PP_ALIGN.CENTER)
txt(s7, "capability routing · fallback · model-agnostic",
    4.9, 3.32, 3.5, 0.28, size=8, color=SUBTLE, align=PP_ALIGN.CENTER)
txt(s7, "worker_manager.py · LiteLLM", 4.9, 3.62, 3.5, 0.22, size=7, color=MUTED, align=PP_ALIGN.CENTER)

# Provider cards (left)
providers = [
    ("Gemini 2.5 Flash",   "google",      "conversational, fast",    ONLINE),
    ("Gemini Flash Latest","google",      "research, analysis",       ONLINE),
    ("Llama 3.3 70B",      "groq",        "fast, open-source",        ONLINE),
    ("Llama 3.2 3B",       "openrouter",  "free, public endpoint",    ONLINE),
    ("Ollama (local)",     "local",       "private, offline",         AMBER),
]
for i, (name, provider, cap, color) in enumerate(providers):
    yp = 1.4 + i * 0.9
    box(s7, 0.2, yp, 3.9, 0.75, fill=PANEL, border_color=color)
    txt(s7, name, 0.35, yp + 0.05, 3.5, 0.28, size=11, bold=True, color=color)
    txt(s7, f"{provider}  ·  {cap}", 0.35, yp + 0.35, 3.5, 0.25, size=8, color=SUBTLE)

# Arrow connector middle  
for i in range(5):
    yp = 1.62 + i * 0.9
    txt(s7, "─────────►", 4.1, yp + 0.1, 0.8, 0.25, size=8, color=ION)

# Output cards (right)
outputs = [
    ("→ Conversational response",    GLOW),
    ("→ Mission extraction",          GLOW),
    ("→ Command generation",          GLOW),
    ("→ Review / validation",         GLOW),
    ("→ Future: Analysis agents",     MUTED),
]
for i, (out, color) in enumerate(outputs):
    yp = 1.4 + i * 0.9
    box(s7, 9.1, yp, 4.0, 0.75, fill=PANEL, border_color=color)
    txt(s7, out, 9.25, yp + 0.22, 3.7, 0.3, size=11, color=color)

# Bottom honest callout
box(s7, 0.15, 6.25, 12.9, 0.75, fill=RGBColor(0x20,0x15,0x03), border_color=AMBER)
txt(s7, "⚡  CURRENT REALITY:  Workers = model endpoints routed via LiteLLM.  "
        "Autonomous specialized worker agents with independent execution loops = FUTURE MILESTONE.",
    0.35, 6.38, 12.6, 0.5, size=10, color=AMBER)

txt(s7, "07 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — REVIEW ENGINE
# ══════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
section_header(s8, "AI DOESN'T GET TO DECLARE ITSELF SUCCESSFUL.", "GENERATION ≠ COMPLETION")

# Pipeline (center column)
stages = [
    ("WORKER OUTPUT",               MUTED,  PANEL),
    ("STAGE 1: EXIT CODE CHECK",    ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("exit_code == 0 ?",            SUBTLE, PANEL),
    ("STAGE 2: ARTIFACT CHECK",     ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("Does output.txt exist on disk?",  SUBTLE, PANEL),
    ("STAGE 3: HEURISTIC CHECK",    ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("Output sanity validation",    SUBTLE, PANEL),
    ("STAGE 4: SEMANTIC LLM REVIEW",GLOW,  DARK_VIOLET),
    ("LLM checks output vs success_criteria", SUBTLE, PANEL),
]

cx = 4.2
y_s = 1.5
for i, (label_text, color, bg_c) in enumerate(stages):
    h_s = 0.42 if "STAGE" in label_text or label_text in ("WORKER OUTPUT",) else 0.3
    box(s8, cx, y_s, 4.9, h_s, fill=bg_c, border_color=color)
    sz = 12 if "STAGE" in label_text else 9
    bold = "STAGE" in label_text or label_text == "WORKER OUTPUT"
    txt(s8, label_text, cx + 0.12, y_s + 0.06, 4.7, h_s - 0.1,
        size=sz, bold=bold, color=color)
    if i < len(stages) - 1:
        txt(s8, "▼", cx + 2.3, y_s + h_s, 0.4, 0.25, size=8, color=ION, align=PP_ALIGN.CENTER)
    y_s += h_s + 0.25

# Pass / Retry at bottom
box(s8, 3.0, 6.55, 2.5, 0.52, fill=RGBColor(0x04,0x22,0x12), border_color=ONLINE)
txt(s8, "✓  PASS → DELIVER", 3.05, 6.65, 2.4, 0.32, size=11, bold=True, color=ONLINE, align=PP_ALIGN.CENTER)

box(s8, 7.8, 6.55, 2.5, 0.52, fill=RGBColor(0x22,0x10,0x04), border_color=AMBER)
txt(s8, "🔄  RETRY → RE-PROMPT", 7.85, 6.65, 2.4, 0.32, size=11, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Right side annotation cards
card(s8, 9.3, 1.55, 3.8, 1.5, "WHY THIS MATTERS",
     ["LLMs generate plausible output —",
      "not guaranteed correct output.",
      "The Review Engine is the boundary",
      "between AI assistant and reliable system."],
     title_color=ONLINE, border=ONLINE)

card(s8, 9.3, 3.2, 3.8, 1.5, "RETRY MECHANISM",
     ["Failed output + error message fed",
      "back to ExecutionBrain as context.",
      "LLM attempts a different command.",
      "Max retries enforced."],
     title_color=GLOW, border=ION)

txt(s8, "review_engine.py", 0.3, 7.05, 4.0, 0.3, size=8, color=MUTED, italic=True)
txt(s8, "08 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — THE INTERFACE
# ══════════════════════════════════════════════════════════════════════════════
s9 = add_slide()
section_header(s9, "THE INTERFACE IS A CONTROL SURFACE", "FOR A LIVING MISSION.")

# Left: flow pipeline
flow = [
    ("VOICE INPUT",            "Web Speech API (browser STT)",   SUBTLE, PANEL),
    ("REACT / VITE UI",        "frontend/src/  —  IMPLEMENTED",   ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("POST  /api/chat",        "HadesService.sendMessage()",       GLOW, DARK_VIOLET),
    ("FASTAPI BACKEND",        "main.py  —  IMPLEMENTED",          ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("PARTNER BRAIN",          "Conversation + Mission Lock",       ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("EXECUTION BRAIN",        "Background async task",             AMBER, RGBColor(0x20,0x15,0x03)),
    ("SSE  /api/events",       "Real-time event stream",           ION, DARK_VIOLET),
    ("LIVE UI UPDATES",        "System Activity feed",             ONLINE, RGBColor(0x04,0x1e,0x10)),
    ("HADES VOICE RESPONSE",   "Kokoro-ONNX → Base64 → HTML5 Audio", ONLINE, RGBColor(0x04,0x1e,0x10)),
]

y_f = 1.55
for step_name, detail, color, bg_c in flow:
    box(s9, 0.2, y_f, 5.5, 0.5, fill=bg_c, border_color=color)
    txt(s9, step_name, 0.35, y_f + 0.04, 5.2, 0.25, size=11, bold=True, color=color)
    txt(s9, detail,    0.35, y_f + 0.28, 5.2, 0.2,  size=8,  color=SUBTLE)
    if step_name != flow[-1][0]:
        txt(s9, "▼", 2.7, y_f + 0.5, 0.4, 0.2, size=8, color=ION, align=PP_ALIGN.CENTER)
    y_f += 0.72

# Right: 4 cards
interface_cards = [
    ("REAL-TIME EVENTS", ONLINE, ONLINE,
     ["Server-Sent Events stream every step to the UI.",
      "Task started → capability selected → completed.",
      "User sees Hades working, not just a spinner."]),
    ("BACKEND TTS — Kokoro-ONNX", GLOW, ION,
     ["Local neural TTS on the backend server.",
      "Audio encoded as Base64 in JSON response.",
      "No cloud TTS dependency required."]),
    ("VOICE INPUT", AMBER, AMBER,
     ["Currently: browser window.SpeechRecognition.",
      "Backend Whisper STT — NOT YET IMPLEMENTED.",
      "Future milestone for offline voice input."]),
    ("SESSION PERSISTENCE", ONLINE, ONLINE,
     ["Session ID stored in localStorage.",
      "Conversation history maintained in session.",
      "Survives page refresh — not server restart."]),
]
y_c = 1.55
for title, title_color, border_color, body in interface_cards:
    card(s9, 6.0, y_c, 7.1, 1.12, title, body,
         title_color=title_color, border=border_color, title_size=11, body_size=9)
    y_c += 1.24

txt(s9, "frontend/src/services/HadesService.ts  ·  main.py  ·  voice_manager.py",
    0.3, 7.08, 12.7, 0.28, size=7, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
txt(s9, "09 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — END-TO-END MISSION PROOF
# ══════════════════════════════════════════════════════════════════════════════
s10 = add_slide()
section_header(s10, "ONE SENTENCE → REAL COMPUTER ACTION", "VERIFIED END-TO-END EXECUTION TRACE")

# Mission statement
box(s10, 0.3, 1.52, 12.7, 0.5, fill=RGBColor(0x0c,0x0c,0x22), border_color=GLOW)
txt(s10, ">  \"List all files in the current directory and save them to output.txt.\"",
    0.5, 1.6, 12.5, 0.35, size=12, color=GLOW)

# 13 steps in 2 columns
steps_e2e = [
    ("01", "USER",       "Types or speaks the command",                          SUBTLE),
    ("02", "FRONTEND",   "HadesService.sendMessage → POST /api/chat",            ION),
    ("03", "API",        "main.py receives request, looks up SessionState by UUID", ION),
    ("04", "PARTNER",    "IntentClassifier → classifies as SMALL_TASK",          GLOW),
    ("05", "PARTNER",    "MissionExtractor → populates MissionUnderstanding fields", GLOW),
    ("06", "LOCK",       "UnderstandingEvaluator → AUTHORIZED_EXECUTION",        WHITE),
    ("07", "API",        "asyncio.create_task(execution_brain.process_mission) — HTTP returns immediately", ION),
    ("08", "EXEC",       "ExecutionBrain._generate_plan → 1-task TaskGraph",     AMBER),
    ("09", "WORKER",     "WorkerManager → Gemini → LLM generates: ls -la > output.txt", ONLINE),
    ("10", "TOOL",       "TerminalSkill.execute() → asyncio.create_subprocess_shell", ONLINE),
    ("11", "REVIEW",     "ReviewEngine: exit_code==0 ✓  ·  output.txt exists on disk ✓", ONLINE),
    ("12", "MEMORY",     "memory_manager.add_mission_to_history() → memory.json",SUBTLE),
    ("13", "SSE",        "EventSource emits MISSION_COMPLETED → UI updates, Hades speaks result", ONLINE),
]

col1_steps = steps_e2e[:7]
col2_steps = steps_e2e[7:]

step_h_e = 0.42
y_e = 2.18
for step in col1_steps:
    num, layer, action, color = step
    box(s10, 0.15, y_e, 6.3, step_h_e, fill=PANEL, border_color=color)
    label(s10, num, 0.2, y_e + 0.07, 0.35, 0.27, bg=color, fg=ABYSS, size=8)
    label(s10, layer, 0.62, y_e + 0.07, 1.2, 0.27, bg=PANEL, fg=color, size=7)
    txt(s10, action, 1.9, y_e + 0.08, 4.5, 0.3, size=9, color=TEXT)
    y_e += step_h_e + 0.05

y_e = 2.18
for step in col2_steps:
    num, layer, action, color = step
    box(s10, 6.7, y_e, 6.45, step_h_e, fill=PANEL, border_color=color)
    label(s10, num, 6.75, y_e + 0.07, 0.35, 0.27, bg=color, fg=ABYSS, size=8)
    label(s10, layer, 7.17, y_e + 0.07, 1.2, 0.27, bg=PANEL, fg=color, size=7)
    txt(s10, action, 8.45, y_e + 0.08, 4.6, 0.3, size=9, color=TEXT)
    y_e += step_h_e + 0.05

# Bottom chain
hline(s10, 0.3, 6.95, 12.7, color=ION, thickness=Pt(0.75))
chain = ["CONVERSATION", "→", "MISSION LOCK", "→", "ASYNC EXECUTION", "→", "VERIFICATION", "→", "DELIVERY"]
chain_colors = [SUBTLE, ION, WHITE, ION, AMBER, ION, ONLINE, ION, ONLINE]
xc = 0.3
for c_item, c_color in zip(chain, chain_colors):
    w_c = len(c_item) * 0.115 + 0.15
    txt(s10, c_item, xc, 7.05, w_c, 0.3, size=10, bold=(c_item != "→"),
        color=c_color, align=PP_ALIGN.CENTER)
    xc += w_c

txt(s10, "10 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MATURITY MATRIX
# ══════════════════════════════════════════════════════════════════════════════
s11 = add_slide()
section_header(s11, "THE PROTOTYPE IS REAL.", "BUT IT IS NOT FINISHED.")

green_items = [
    "Conversation Flow", "Intent Classification", "Mission Extraction",
    "Mission State Machine (5 states)", "Mission Lock Gate",
    "WorkerManager + LiteLLM", "Multi-Provider Routing (Gemini, Groq, OpenRouter, Ollama)",
    "Terminal Skill", "Filesystem Skill", "Process Manager",
    "Review Engine (4 stages)", "Server-Sent Events stream",
    "React / Vite Frontend", "Full E2E execution path",
    "Backend TTS — Kokoro-ONNX", "Session persistence (localStorage)",
]
amber_items = [
    "Worker specialization (endpoints only; no agent loops)",
    "Task decomposition (hardcoded to 1 task)",
    "Session memory (in-memory; lost on restart)",
    "Mission memory (flat JSON; no semantic retrieval)",
    "Voice input (browser SpeechRecognition only)",
    "Browser skill (Playwright wrapper, unverified E2E)",
    "Equal partnership guardrails (prompt-based only)",
]
red_items = [
    "Dynamic multi-step TaskGraph decomposition",
    "Autonomous specialized worker agents",
    "RAG / Vector semantic memory",
    "Backend Whisper STT",
    "Docker / sandbox isolation",
    "Execution security hardening",
    "Deep Linux / systemd / D-Bus integration",
    "Structured logging + observability",
    "Comprehensive test suite (1 file exists)",
    "Production-grade automation",
    "Cross-session persistent identity",
]

col_configs = [
    ("✓  IMPLEMENTED",  green_items,  ONLINE, RGBColor(0x04,0x18,0x0e)),
    ("⚡  PARTIAL",      amber_items,  AMBER,  RGBColor(0x18,0x10,0x02)),
    ("✗  MISSING",      red_items,    DANGER, RGBColor(0x18,0x04,0x04)),
]

col_xs = [0.15, 4.55, 8.95]
for ci, (header, items, color, bg_c) in enumerate(col_configs):
    # Column header
    box(s11, col_xs[ci], 1.52, 4.2, 0.42, fill=bg_c, border_color=color)
    txt(s11, header, col_xs[ci] + 0.1, 1.56, 4.0, 0.34, size=12, bold=True, color=color)
    # Items
    yi = 2.0
    for item in items:
        if yi > 6.8:
            break
        box(s11, col_xs[ci], yi, 4.2, 0.36, fill=bg_c, border_color=RGBColor(0x30,0x30,0x40))
        txt(s11, item, col_xs[ci] + 0.12, yi + 0.05, 4.0, 0.28, size=8.5, color=TEXT)
        yi += 0.38

# Bottom statement
hline(s11, 0.15, 6.95, 12.95, color=MUTED, thickness=Pt(0.5))
label(s11, "CORE LOOP: PROVEN", 3.3, 7.08, 2.5, 0.3,
      bg=RGBColor(0x04,0x18,0x0e), fg=ONLINE, size=10)
label(s11, "AI OS: IN PROGRESS", 7.4, 7.08, 2.5, 0.3,
      bg=DARK_VIOLET, fg=GLOW, size=10)

txt(s11, "11 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — NOVELTY CARDS
# ══════════════════════════════════════════════════════════════════════════════
s12 = add_slide()
section_header(s12, "WHAT HADES ACTUALLY CONTRIBUTES.", "7 ARCHITECTURAL NOVELTIES")

novelties = [
    ("01", "🔒", "MUTUAL UNDERSTANDING GATE",
     "Conversation is the authorization layer. Execution is gated by a deterministic "
     "state check: objective + desired_outcome + success_criteria + mutual_understanding. "
     "No other AI system gates tool use on explicit human-alignment verification.",
     "IMPLEMENTED", ONLINE),
    ("02", "🎯", "MISSION-CENTRIC COMPUTING",
     "Prompts are disposable. Missions are persistent units with structured fields: "
     "objective, constraints, success criteria, state, and outcome. The mission survives "
     "the full conversation → execution → review → memory lifecycle.",
     "IMPLEMENTED", ONLINE),
    ("03", "🧠", "PARTNER + EXECUTIVE DUALITY",
     "Two isolated cognitive systems. Partner Brain handles conversation — never touches tools. "
     "Executive Brain handles planning — never engages user directly. Prevents the classic "
     "agent failure: execution and conversation contaminating each other.",
     "IMPLEMENTED", ONLINE),
    ("04", "⚙️", "INVISIBLE ORCHESTRATION",
     "User expresses a goal in natural language. Hades internally classifies capability, "
     "selects the right model, routes the prompt, collects results, validates output — "
     "without the user knowing which model was used. Intelligence becomes infrastructure.",
     "IMPLEMENTED", ONLINE),
    ("05", "✅", "VERIFICATION-FIRST EXECUTION",
     "Workers generate output. Review Engine determines if the MISSION was completed. "
     "4-stage validation: exit code → artifact existence → heuristic → LLM semantic review. "
     "On failure, system retries with error as context. Generation ≠ Success.",
     "IMPLEMENTED", ONLINE),
    ("06", "🔄", "MODEL-AGNOSTIC AI LAYER",
     "LiteLLM wraps Google, Groq, OpenRouter, and local Ollama behind a unified interface. "
     "Fallback is automatic. New models added via config.json without code changes. "
     "The system is not architecturally coupled to any vendor.",
     "IMPLEMENTED", ONLINE),
    ("07", "🤝", "PERSISTENT RELATIONSHIP",
     "Mission history, conversation context, and user identity persist across interactions. "
     "Hades addresses users by name and recalls past missions. Lays the foundation for "
     "genuine long-term AI partnership — not session-scoped Q&A. (RAG = future milestone)",
     "PARTIAL", AMBER),
]

# 4 + 3 layout
positions = [
    (0.15, 1.52, 3.2, 1.75),
    (3.48, 1.52, 3.2, 1.75),
    (6.81, 1.52, 3.2, 1.75),
    (10.14, 1.52, 3.1, 1.75),
    (0.15, 3.38, 4.15, 1.75),
    (4.43, 3.38, 4.15, 1.75),
    (8.71, 3.38, 4.55, 1.75),
]

for i, (nov, pos) in enumerate(zip(novelties, positions)):
    num, icon, title, body, status, color = nov
    x_c, y_c, w_c, h_c = pos
    border_c = color
    bg_c = RGBColor(0x04,0x18,0x0e) if color == ONLINE else RGBColor(0x18,0x10,0x02)
    box(s12, x_c, y_c, w_c, h_c, fill=PANEL, border_color=border_c)
    txt(s12, f"{num}  {icon}", x_c + 0.1, y_c + 0.07, 0.8, 0.3, size=11, bold=True, color=MUTED)
    txt(s12, title, x_c + 0.1, y_c + 0.38, w_c - 0.2, 0.38, size=10, bold=True, color=color)
    txt(s12, body, x_c + 0.1, y_c + 0.78, w_c - 0.2, h_c - 0.95, size=7.5, color=SUBTLE)

# Bottom
hline(s12, 0.3, 5.25, 12.7, color=MUTED, thickness=Pt(0.5))
txt(s12, "The novelty is not another smarter model.",
    0.3, 5.32, 12.7, 0.38, size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
txt(s12, "It is a system that organizes intelligence around human intent.",
    0.3, 5.72, 12.7, 0.38, size=16, bold=True, color=GLOW, align=PP_ALIGN.CENTER)

txt(s12, "12 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ROADMAP + CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s13 = add_slide()
section_header(s13, "FROM WORKING PROTOTYPE", "TO TRUE AI OPERATING SYSTEM")

cols_r = [
    ("v0.1  ·  TODAY", ONLINE, RGBColor(0x04,0x1e,0x10),
     ["✓  Partner Brain conversational loop",
      "✓  Mission Lock authorization gate",
      "✓  Multi-provider model routing",
      "✓  Terminal + filesystem execution",
      "✓  Review Engine (4-stage validation)",
      "✓  Real-time SSE event stream",
      "✓  Backend TTS — Kokoro-ONNX",
      "✓  Complete E2E execution path"],
     "CORE LOOP: PROVEN"),
    ("v0.2  ·  NEXT", AMBER, RGBColor(0x20,0x15,0x03),
     ["→  Dynamic TaskGraph decomposition",
      "→  True multi-step mission execution",
      "→  Docker sandbox (safe terminal)",
      "→  Whisper backend STT",
      "→  Specialized autonomous workers",
      "→  Vector / RAG semantic memory",
      "→  Security hardening",
      "→  E2E test suite"],
     "KNOWN ENGINEERING PROBLEMS"),
    ("v1.0  ·  AI OS VISION", GLOW, DARK_VIOLET,
     ["◆  Persistent cross-device identity",
      "◆  Autonomous mission scheduling",
      "◆  Event-driven background missions",
      "◆  Plugin / skill ecosystem",
      "◆  Distributed multi-node execution",
      "◆  Local + cloud intelligence blend",
      "◆  Long-term human-AI partnership",
      "◆  Deep Linux / systemd integration"],
     "THE ACTUAL AI OS VISION"),
]

col_w_r = 4.2
for i, (title, color, bg_c, items, badge_text) in enumerate(cols_r):
    x_r = 0.15 + i * (col_w_r + 0.12)
    # Header
    box(s13, x_r, 1.52, col_w_r, 0.5, fill=bg_c, border_color=color)
    txt(s13, title, x_r + 0.1, 1.57, col_w_r - 0.2, 0.4, size=13, bold=True, color=color)
    # Items
    y_r = 2.1
    for item in items:
        box(s13, x_r, y_r, col_w_r, 0.42, fill=PANEL, border_color=RGBColor(0x28,0x28,0x40))
        txt(s13, item, x_r + 0.1, y_r + 0.08, col_w_r - 0.2, 0.28, size=10, color=TEXT)
        y_r += 0.44
    # Badge
    label(s13, badge_text, x_r, y_r + 0.1, col_w_r, 0.3, bg=bg_c, fg=color, size=8)

# Final statement
hline(s13, 0.3, 6.3, 12.7, color=ION, thickness=Pt(1.5))
txt(s13, "HUMAN SETS THE MISSION.",
    0.3, 6.42, 12.7, 0.52, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s13, "HADES ORCHESTRATES THE INTELLIGENCE.",
    0.3, 6.9, 12.7, 0.52, size=24, bold=True, color=GLOW, align=PP_ALIGN.CENTER)

hline(s13, 0.3, 7.38, 12.7, color=ION, thickness=Pt(1.5))

txt(s13, "13 / 13", 12.5, 7.1, 0.7, 0.3, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ── SAVE ─────────────────────────────────────────────────────────────────────
output_path = r"d:\HACK O HADES\HADES_PRESENTATION.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
